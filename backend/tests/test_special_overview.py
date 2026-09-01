"""专项总览：七列全部自动提取，只有点灯能人工覆盖。

这张表的价值全在「自动」上——一旦哪一列变成要人另填一遍，总览就会
和详情页对不上，而两边看着都对。所以这里断言的重点是：

1. 六列文本确实来自各专项自己的字段（并且富文本被剥成纯文本）；
2. 点灯的四档规则（未评估 / 绿 / 黄 / 红）以及两条容易被"顺手改掉"的边界：
   当天到期不算超期、没填计划闭环时间的未闭环行只到黄；
3. 手工覆盖优先、清空能回到自动，并且走 content 的乐观锁；
4. `/overview` 注册在 `/{sid}` 前面——顺序反了这个接口会稳定 422，
   而错误信息只会说"sid 不是整数"，没人会往路由顺序上想。
"""
from datetime import date, timedelta

import pytest


def _mk_special(client, admin_headers, name, **kw):
    body = {"name": name, "kind": "special", "owner": "张三"}
    body.update(kw)
    r = client.post("/api/specials", json=body, headers=admin_headers)
    assert r.status_code == 200, r.text
    return r.json()["id"]


def _add_risk(client, admin_headers, sid, content, **kw):
    body = {"content": content}
    body.update(kw)
    r = client.post(f"/api/specials/{sid}/risks", json=body, headers=admin_headers)
    assert r.status_code == 200, r.text
    return r.json()


def _overview(client, admin_headers, include_inactive=False):
    r = client.get("/api/specials/overview",
                   params={"include_inactive": include_inactive},
                   headers=admin_headers)
    assert r.status_code == 200, r.text
    return {row["id"]: row for row in r.json()}


def _row(client, admin_headers, sid):
    rows = _overview(client, admin_headers)
    assert sid in rows, f"{sid} 应出现在总览里：{list(rows)}"
    return rows[sid]


def _set_content(client, admin_headers, sid, **fields):
    cur = client.get(f"/api/specials/{sid}", headers=admin_headers).json()["content"]
    body = {"version": cur["version"], **fields}
    r = client.put(f"/api/specials/{sid}/content", json=body, headers=admin_headers)
    assert r.status_code == 200, r.text
    return r.json()


@pytest.fixture(scope="module")
def sid(client, admin_headers):
    return _mk_special(client, admin_headers, "总览用例专项")


# ─── 路由顺序 ────────────────────────────────────────────────────────────


def test_overview_is_not_parsed_as_sid(client, admin_headers):
    """`/overview` 必须在 `/{sid}` 之前注册，否则稳定 422。"""
    r = client.get("/api/specials/overview", headers=admin_headers)
    assert r.status_code == 200, r.text
    assert isinstance(r.json(), list)


def test_overview_pptx_is_not_parsed_as_sid(client, admin_headers):
    """`/overview.pptx` 同样必须排在 `/{sid}` 之前。"""
    r = client.get("/api/specials/overview.pptx", headers=admin_headers)
    assert r.status_code == 200, r.text
    assert r.content[:2] == b"PK", "返回的应该是个 pptx（zip）"
    assert "attachment" in r.headers.get("content-disposition", "")


def test_overview_pptx_uses_the_same_rows_as_the_page(client, admin_headers):
    """导出的行数与页面一致，停用的同样默认不进。

    导出里另查一遍库的话，两处的过滤迟早分叉——页面上 5 行、PPT 里 6 行，
    而两边看着都对。
    """
    from pptx import Presentation
    import io as _io

    page_rows = _overview(client, admin_headers)
    r = client.get("/api/specials/overview.pptx", headers=admin_headers)
    pres = Presentation(_io.BytesIO(r.content))
    body = 0
    for slide in pres.slides:
        for sh in slide.shapes:
            if sh.has_table:
                body += len(sh.table.rows) - 2      # 两行表头
    assert body == len(page_rows)


# ─── 六列的取数 ──────────────────────────────────────────────────────────


def test_columns_come_from_the_special_itself(client, admin_headers, sid):
    _set_content(client, admin_headers, sid,
                 goal="<p>把 A 模块的时延压到 <b>10ms</b> 以内</p>",
                 progress_summary="<p>已完成一轮压测</p>")
    row = _row(client, admin_headers, sid)
    assert row["name"] == "总览用例专项"
    assert row["owner"] == "张三"
    assert row["kind_label"] == "专项"
    # 富文本剥成纯文本：这张表是横着扫的，标签留在里面每行高矮不一
    assert row["goal"] == "把 A 模块的时延压到 10ms 以内"
    assert "<b>" not in row["goal"]
    assert row["progress"] == "已完成一轮压测"


def test_seq_follows_sidebar_order(client, admin_headers):
    """序号＝表内位置（1 起、连续），不是 id——总览与侧栏对不上人就得两边找。"""
    r = client.get("/api/specials/overview", headers=admin_headers)
    seqs = [x["seq"] for x in r.json()]
    assert seqs == list(range(1, len(seqs) + 1))


# ─── 点灯的四档 ──────────────────────────────────────────────────────────


def test_no_risk_row_is_gray_not_green(client, admin_headers):
    """一条风险都没登记＝**未评估**，不是绿。

    记成绿的话，最该被追着去补风险的那几个专项在总览上看着比谁都干净。
    """
    s = _mk_special(client, admin_headers, "灯-未登记")
    row = _row(client, admin_headers, s)
    assert row["light"] == "gray"
    assert row["risk_total"] == 0


def test_all_closed_is_green(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-全闭环")
    _add_risk(client, admin_headers, s, "已解决的风险", status="closed")
    row = _row(client, admin_headers, s)
    assert row["light"] == "green"
    assert row["risk_open"] == 0
    # 已闭环的行不进「关键风险和措施」这一列：那一列问的是"现在还有什么风险"
    assert row["risks"] == []


def test_open_but_not_due_is_yellow(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-未到期")
    future = (date.today() + timedelta(days=10)).strftime("%Y-%m-%d")
    _add_risk(client, admin_headers, s, "还在处理", planned_close_date=future)
    row = _row(client, admin_headers, s)
    assert row["light"] == "yellow"
    assert row["risk_open"] == 1 and row["risk_overdue"] == 0


def test_overdue_open_is_red(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-已超期")
    past = (date.today() - timedelta(days=3)).strftime("%Y-%m-%d")
    _add_risk(client, admin_headers, s, "拖了三天", planned_close_date=past)
    row = _row(client, admin_headers, s)
    assert row["light"] == "red"
    assert row["risk_overdue"] == 1
    assert row["risks"][0]["overdue"] is True


def test_due_today_is_not_overdue(client, admin_headers):
    """「超过」才算，当天到期不算——记成超期会让人白紧张一天。"""
    s = _mk_special(client, admin_headers, "灯-今天到期")
    today = date.today().strftime("%Y-%m-%d")
    _add_risk(client, admin_headers, s, "今天到期", planned_close_date=today)
    row = _row(client, admin_headers, s)
    assert row["light"] == "yellow", "当天到期不该记成超期"


def test_missing_plan_date_stops_at_yellow(client, admin_headers):
    """没填计划闭环时间的未闭环行只顶到黄。

    那是「没排期」不是「排了没做到」。混成一档会让一批从来不填日期的专项
    长期挂红，红灯一多就没人看了，真红的那个反而淹掉。
    """
    s = _mk_special(client, admin_headers, "灯-没填日期")
    _add_risk(client, admin_headers, s, "没排期的风险")
    row = _row(client, admin_headers, s)
    assert row["light"] == "yellow"
    assert row["risk_overdue"] == 0


def test_overdue_risks_sort_first(client, admin_headers):
    """一格里放不下几条，先露出来的得是最该看的那几条。"""
    s = _mk_special(client, admin_headers, "灯-排序")
    _add_risk(client, admin_headers, s, "不急的")
    past = (date.today() - timedelta(days=1)).strftime("%Y-%m-%d")
    _add_risk(client, admin_headers, s, "超期的", planned_close_date=past)
    row = _row(client, admin_headers, s)
    assert row["risks"][0]["content"] == "超期的"


def test_measure_column_keeps_risk_and_progress_apart(client, admin_headers):
    """风险与措施是两列，合成一段字符串前端就再也拆不开了。"""
    s = _mk_special(client, admin_headers, "灯-措施")
    _add_risk(client, admin_headers, s, "供货周期长",
              progress="<p>已改用备选料</p>", owner="李四")
    r = _row(client, admin_headers, s)["risks"][0]
    assert r["content"] == "供货周期长"
    assert r["progress"] == "已改用备选料"      # 富文本同样剥成纯文本
    assert r["owner"] == "李四"


# ─── 手工覆盖 ────────────────────────────────────────────────────────────


def test_manual_light_overrides_auto_and_can_be_cleared(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-覆盖")
    _add_risk(client, admin_headers, s, "自动会算成黄")
    row = _row(client, admin_headers, s)
    assert row["light"] == "yellow" and row["light_manual"] == ""

    r = client.put(f"/api/specials/{s}/overview",
                   json={"light": "red", "version": row["version"]},
                   headers=admin_headers)
    assert r.status_code == 200, r.text
    out = r.json()
    assert out["light"] == "red"
    assert out["light_manual"] == "red"
    # 自动档要一并给出去，页面才说得清"这灯是人拍的还是算出来的"
    assert out["light_auto"] == "yellow"

    # 清空＝回到自动
    r = client.put(f"/api/specials/{s}/overview",
                   json={"light": "", "version": out["version"]},
                   headers=admin_headers)
    assert r.status_code == 200, r.text
    assert r.json()["light"] == "yellow"
    assert r.json()["light_manual"] == ""


def test_manual_light_uses_optimistic_lock(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-乐观锁")
    row = _row(client, admin_headers, s)
    ok = client.put(f"/api/specials/{s}/overview",
                    json={"light": "green", "version": row["version"]},
                    headers=admin_headers)
    assert ok.status_code == 200
    stale = client.put(f"/api/specials/{s}/overview",
                       json={"light": "red", "version": row["version"]},
                       headers=admin_headers)
    assert stale.status_code == 409


def test_manual_light_rejects_unknown_value(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-非法值")
    row = _row(client, admin_headers, s)
    r = client.put(f"/api/specials/{s}/overview",
                   json={"light": "紫", "version": row["version"]},
                   headers=admin_headers)
    assert r.status_code == 400


def test_manual_light_accepts_chinese_spelling(client, admin_headers):
    """接口直填「红」也认——前端只传 key，这条是给导入/脚本留的余地。"""
    s = _mk_special(client, admin_headers, "灯-中文写法")
    row = _row(client, admin_headers, s)
    r = client.put(f"/api/specials/{s}/overview",
                   json={"light": "红", "version": row["version"]},
                   headers=admin_headers)
    assert r.status_code == 200, r.text
    assert r.json()["light_manual"] == "red"


# ─── 停用 ────────────────────────────────────────────────────────────────


def test_inactive_specials_hidden_by_default(client, admin_headers):
    s = _mk_special(client, admin_headers, "灯-停用的", is_active=False)
    assert s not in _overview(client, admin_headers)
    assert s in _overview(client, admin_headers, include_inactive=True)
