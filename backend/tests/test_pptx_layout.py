"""PPT 导出的版面：**每一页都要装得下**。

以前所有数据行都画在一张 slide 上，行一多表格就顺着页面往下长、长到幻灯片外面去：
导出来的 PPT 一半内容看不见，得手工拆页——"格式没办法直接用"说的就是这个。
列宽也是手写的英寸数，加一列就对不上，表格要么越过右边界要么右边空一条。

这里钉住的是：
1. 每一页的表格底边不越过幻灯片底边、右边不越过页宽（**行多、单元格文字多都不行**）；
2. 列宽归一化成正好铺满页宽，改列不用重算英寸；
3. 长文本按行截断并写明"另 N 条"，一格吃不掉一整页；
4. 一条数据都没有时给一张写着"暂无数据"的页，而不是一张只有表头的空表；
5. 宽矩阵（客户/月份很多）按列切页，每页都带标签列，列不会被压到看不清；
6. 合计行加粗且有底色，不会混在斑马纹里认不出来。
"""
from types import SimpleNamespace as NS

import pytest
from pptx import Presentation

import pptx_utils as PU

_EMU_IN = 914400
_SLIDE_H_IN = 7.5
_SLIDE_W_IN = 13.334


def _tables(pres):
    """[(页号, table shape)]，封面之类没有表的页跳过。"""
    out = []
    for n, slide in enumerate(pres.slides, 1):
        for shape in slide.shapes:
            if shape.has_table:
                out.append((n, shape))
    return out


def _assert_fits(pres):
    """每一页的表都在幻灯片里。行高用**显式写入的高度**求和——

    PowerPoint 只会把行撑得更高、不会压缩，所以显式高度是下界：
    这个下界都超了，实际导出必然溢出。
    """
    tables = _tables(pres)
    assert tables, "至少应有一张表"
    for page, shape in tables:
        bottom = (shape.top + sum(r.height for r in shape.table.rows)) / _EMU_IN
        right = (shape.left + sum(c.width for c in shape.table.columns)) / _EMU_IN
        assert bottom <= _SLIDE_H_IN, f"第 {page} 页表格越过幻灯片底边（{bottom:.2f}\"）"
        assert right <= _SLIDE_W_IN, f"第 {page} 页表格越过幻灯片右边（{right:.2f}\"）"


def _machine(i, n_issues=0):
    issues = [NS(kind="task", status="OPEN", sort_order=k, id=k,
                 description=f"事务 {i}-{k}：需要配合客户完成联调验证并回归全部用例")
              for k in range(n_issues)]
    return NS(machine_id=f"M{i:03d}", battlefield="某某客户面", model="YLS3000",
              current_stage="量产爬坡", field_version="C10SPC101B003", attention_level=3,
              customer_status="本周完成三个模块联调，下周进入稳定性测试。", issues=issues,
              issue_url="https://example.com/x")


# ─── 客户面状态 ─────────────────────────────────────────────────────────────
def test_customer_status_paginates_instead_of_overflowing():
    pres = Presentation(PU.build_customer_status_pptx([_machine(i) for i in range(40)]))
    assert len(pres.slides) > 1, "40 台机器必须分页，一页画不下"
    _assert_fits(pres)


def test_one_fat_cell_cannot_eat_a_whole_page():
    """一台机器挂 40 条事务：截断到固定行数，否则那一行能顶满整张幻灯片。"""
    pres = Presentation(PU.build_customer_status_pptx(
        [_machine(0, n_issues=40)] + [_machine(i) for i in range(1, 6)]))
    _assert_fits(pres)
    cell = _tables(pres)[0][1].table.cell(2, 7)     # 第一行数据的「现场关键事务」
    assert cell.text.count("\n") + 1 <= PU._DEFAULT_MAX_LINES
    assert "另" in cell.text and "条" in cell.text, "截断了就要写明还剩多少，别悄悄少几行"


def test_empty_export_says_so():
    pres = Presentation(PU.build_customer_status_pptx([]))
    assert not _tables(pres), "没有数据就不该画一张只有表头的空表"
    texts = " ".join(sh.text_frame.text for sl in pres.slides for sh in sl.shapes
                     if sh.has_text_frame)
    assert "没有可导出的机台" in texts


# ─── 迭代需求 ───────────────────────────────────────────────────────────────
def _dreq(i):
    return NS(seq=i, req_no=f"DR{i:04d}", owner="李四", owner_group="PL甲组",
              title=f"领域需求 {i}：某模块支持新的配置下发通道并兼容老版本协议",
              priority="P1", planned_version="C10SPC101B002",
              progress_walkthrough="已完成", progress_reverse="进行中", progress_stc="未开始",
              progress_coding="已延期", progress_bbit="不涉及", progress_clarify="已变更",
              remark="存在变更：客户侧接口未定")


def _preq(i):
    return NS(seq=i, req_no=f"PR{i:04d}", title=f"产品需求 {i}", planned_version="C10SPC101",
              priority="P0", feature="特性A", feature_fo="王五", feature_se="赵六",
              feature_tfo="孙七", code_areas="驱动/协议栈", key_risks="第三方库版本未锁定",
              progress_walkthrough="已完成", progress_reverse="已完成", progress_domain="进行中",
              progress_coding="进行中", progress_joint_debug="未开始", progress_clarify="未开始",
              progress_test_result="未开始", estimated_loc="3000", actual_loc="2750",
              actual_effort="12人天")


def test_iteration_export_fits_every_page():
    it = NS(year=2026, month=3, name="三月迭代", owner="张三")
    pres = Presentation(PU.build_iteration_pptx(
        it, [_dreq(i) for i in range(1, 45)], [_preq(i) for i in range(1, 20)]))
    _assert_fits(pres)
    assert len(pres.slides) > 3


def test_progress_group_header_covers_only_progress_columns():
    """产品需求的「交付进展跟踪」原来把代码量/工作量三列也罩了进去——
    父表头罩错列，看的人会以为那三列也是进展状态。"""
    it = NS(year=2026, month=3, name="", owner="")
    pres = Presentation(PU.build_iteration_pptx(it, [_dreq(1)], [_preq(1)]))
    progress_slide = [s for _, s in _tables(pres)][-1].table
    merged = progress_slide.cell(0, 3)
    assert merged.text == "交付进展跟踪"
    assert progress_slide.cell(0, 10).text == "预估代码量", "代码量列不属于交付进展跟踪"


# ─── 列宽与宽矩阵 ───────────────────────────────────────────────────────────
def test_column_widths_are_normalised_to_page_width():
    """列宽是比例，内部归一化——手写英寸数的话改一列就再也铺不满/或越界。"""
    got = PU._norm_widths([3, 1, 1], 3)
    assert sum(got) == pytest.approx(PU._TABLE_W_IN)
    assert got[0] == pytest.approx(got[1] * 3)
    # 数量对不上时退化成等宽，而不是照着错的用
    assert PU._norm_widths([1, 2], 3) == pytest.approx([PU._TABLE_W_IN / 3] * 3)


def test_wide_matrix_splits_by_columns_and_keeps_label():
    pres = PU._new_pres()
    cols = [f"客户面{i}" for i in range(30)]
    rows = [[f"PL{g}组"] + [str(g * i) for i in range(30)] for g in range(8)]
    rows.append(["合计"] + ["99"] * 30)
    PU.add_matrix_slides(pres, "按客户分布", "测试", "小组", cols, rows)
    _assert_fits(pres)
    tables = _tables(pres)
    assert len(tables) > 1, "30 列必须按列切页，挤在一页每列只剩 0.4 英寸"
    for _, shape in tables:
        assert shape.table.cell(0, 0).text == "小组", "每一页都要带上标签列，否则数字没有归属"
        assert shape.table.columns[1].width / _EMU_IN >= PU._MIN_MATRIX_COL_IN


def test_total_row_is_visually_distinct():
    pres = PU._new_pres()
    PU.add_matrix_slides(pres, "t", "s", "小组", ["A", "B"],
                         [["PL1组", "1", "2"], ["合计", "1", "2"]])
    table = _tables(pres)[0][1].table
    total_run = table.cell(3, 0).text_frame.paragraphs[0].runs[0]
    assert total_run.font.bold, "合计行混在斑马纹里认不出来"


# ─── 行高估算 ───────────────────────────────────────────────────────────────
def test_line_estimate_counts_cjk_as_full_width():
    """中文按整宽算。按 ASCII 宽度估的话会算出"装得下"，导出的表照样溢出。"""
    assert PU._wrapped_lines("中" * 20, 1.0, 10) > PU._wrapped_lines("a" * 20, 1.0, 10)
    assert PU._wrapped_lines("a\nb\nc", 5.0, 10) == 3


def test_url_column_becomes_a_short_hyperlink():
    """问题单列原来把整条 URL 摊在格子里：既把列撑宽，又没人会照着念。"""
    pres = Presentation(PU.build_customer_status_pptx([_machine(0)]))
    cell = _tables(pres)[0][1].table.cell(2, 9)
    assert cell.text == "查看"
    run = cell.text_frame.paragraphs[0].runs[0]
    assert run.hyperlink.address == "https://example.com/x"


def test_missing_url_still_renders_placeholder():
    m = _machine(1)
    m.issue_url = ""
    pres = Presentation(PU.build_customer_status_pptx([m]))
    assert _tables(pres)[0][1].table.cell(2, 9).text == "—"


# ─── 配色：对齐部门述职模板 ─────────────────────────────────────────────────
def _solid_rgb(shape):
    """形状的纯色填充色（6 位十六进制）；没有填充或不是纯色返回 None。"""
    try:
        return str(shape.fill.fore_color.rgb)
    except (TypeError, AttributeError, ValueError):
        return None


def test_status_is_a_cell_fill_not_just_coloured_text():
    """进展状态**给格子上底色**。

    只给字上色的话，6 个进展列全是同一个字号的小字，得逐格去读才分得出来
    ——而这一屏最该一眼看出的就是哪几格是绿的。这个断言防的是
    "有人觉得底色太花，改回染字色"：改完导出的 PPT 一样能打开、一样有内容。
    """
    it = NS(year=2026, month=3, name="三月迭代", owner="张三")
    pres = Presentation(PU.build_iteration_pptx(it, [_dreq(1)]))
    table = _tables(pres)[0][1].table
    fills = {}
    for col in range(7, 13):
        cell = table.cell(2, col)
        fills[cell.text] = str(cell.fill.fore_color.rgb)

    assert fills["已完成"] == "92D050", "已完成＝绿"
    assert fills["进行中"] == "FFD966", "进行中＝黄"
    assert fills["已延期"] == "FF9999", "已延期＝红"
    # 「已变更」＝本轮不做了。导出不剔这些行（那是交付记录），靠灰底把它和
    # 还在推进的行区分开，否则它混在里面会被当成进度算进去。
    assert fills["已变更"] == "D9D9D9", "已变更＝灰"
    # 这两档表达的是"这里没有进展"，上了底色会跟真有状态的格子一样抢眼
    assert fills["未开始"] in ("FFFFFF", str(PU._ZEBRA)), "未开始不点灯"
    assert fills["不涉及"] in ("FFFFFF", str(PU._ZEBRA)), "不涉及不点灯"


# ─── 专项总览 ────────────────────────────────────────────────────────────


def _orisk(text, act="", overdue=False):
    return NS(id=1, content=text, progress=act, owner="李工",
              planned_close_date="2026-08-28", overdue=overdue)


def _orow(seq, light="yellow", risks=(), risk_total=None, manual="", kind="special"):
    risks = list(risks)
    return NS(seq=seq, id=seq, name=f"专项 {seq}", kind=kind,
              kind_label="攻关" if kind == "assault" else "专项",
              owner="张三",
              goal="把一次通过率从 82% 提到 95% 以上，9 月底前收口。",
              progress="已完成三轮 DOE，本周起在 2 号线试产验证。",
              risks=risks, light=light, light_auto=light, light_manual=manual,
              light_reason="", risk_total=len(risks) if risk_total is None else risk_total,
              risk_open=len(risks), risk_overdue=0, version=0)


def test_special_overview_fits_every_page():
    """行多、每格文字也多——照样每页都装得下。"""
    rows = [_orow(i, risks=[_orisk(f"风险 {i}-{k}：供货周期长且备料只够两轮试产",
                                   "已提前下单并启动并行询价", overdue=(k == 0))
                            for k in range(6)])
            for i in range(1, 26)]
    pres = Presentation(PU.build_special_overview_pptx(rows))
    _assert_fits(pres)
    assert len(_tables(pres)) > 1, "25 个专项该分页，而不是挤在一张里"


def test_special_overview_light_is_a_cell_fill():
    """风险那一列**给格子上底色**，四档各一色。

    这套灯与进展状态词的点灯（`_STATUS_FILLS`）是两回事，配色在
    `brand.LIGHT_FILLS`——合并的话，改一个状态词的颜色会顺手改掉这里。
    """
    rows = [_orow(1, light="red"), _orow(2, light="yellow"),
            _orow(3, light="green"), _orow(4, light="gray")]
    pres = Presentation(PU.build_special_overview_pptx(rows))
    table = _tables(pres)[0][1].table
    got = {}
    for i in range(len(rows)):
        cell = table.cell(2 + i, 3)          # 第 4 列＝风险
        got[cell.text] = str(cell.fill.fore_color.rgb)
    assert got == {"红": "FEF0F0", "黄": "FDF6EC",
                   "绿": "F0F9EB", "未评估": "F4F4F5"}


def test_special_overview_gray_reads_as_not_assessed():
    """一条风险都没登记的那行写「未评估」，不是绿、也不是空白。"""
    pres = Presentation(PU.build_special_overview_pptx([_orow(1, light="gray")]))
    table = _tables(pres)[0][1].table
    assert table.cell(2, 3).text == "未评估"
    assert table.cell(2, 5).text == "未登记风险"


def test_special_overview_says_how_many_risks_were_clipped():
    """一格里塞不下的风险要写明「另 N 条」，而且那个数按**条**算不是按行算。

    风险与措施接在同一行里就是为了这个：一条占两行的话，截断时会写成
    "另 4 条"而其实只剩 2 条，那个数看着完全合理，没人会去核。
    """
    rows = [_orow(1, risks=[_orisk(f"风险 {k}：这是一条足够长的风险描述，"
                                   f"长到一行放不下需要折行显示才够", "对应的措施也写得挺长")
                            for k in range(12)])]
    pres = Presentation(PU.build_special_overview_pptx(rows))
    text = _tables(pres)[0][1].table.cell(2, 5).text
    assert "另" in text and "条" in text, f"截断了却没说少了几条：{text!r}"
    shown = sum(1 for ln in text.split("\n") if ln and ln[0].isdigit())
    rest = int(text.split("另 ")[1].split(" 条")[0])
    assert shown + rest == 12, f"露出 {shown} 条 + 另 {rest} 条 ≠ 12 条"


def test_special_overview_reports_manual_lights_in_the_subtitle():
    """有人工拨过的灯就要说一句。

    不说的话，这份材料把"人拍的板"摆成了"系统算出来的结论"，
    而看的人分不出哪几行是哪种。
    """
    def _sub(pres):
        texts = []
        for shape in pres.slides[0].shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        return "\n".join(texts)

    plain = Presentation(PU.build_special_overview_pptx([_orow(1)]))
    assert "人工指定" not in _sub(plain), "一个都没人工拨过就别提这句"
    mixed = Presentation(PU.build_special_overview_pptx(
        [_orow(1), _orow(2, light="red", manual="red")]))
    assert "其中 1 项的灯为人工指定" in _sub(mixed)


def test_special_overview_empty_says_so():
    pres = Presentation(PU.build_special_overview_pptx([]))
    assert len(pres.slides) == 1
    body = "\n".join(sh.text_frame.text for sh in pres.slides[0].shapes
                     if sh.has_text_frame)
    assert "还没有专项" in body


def test_header_is_light_blue_with_dark_text_not_red_on_white():
    """一页里的红只留给标题。表头一红，整张表的重心就压在最上面一行。"""
    pres = Presentation(PU.build_customer_status_pptx([_machine(0)]))
    # 客户面表没有父表头，首列是 0/1 行纵向合并的，文字在 cell(0,0)
    cell = _tables(pres)[0][1].table.cell(0, 0)
    assert str(cell.fill.fore_color.rgb) == str(PU._HEADER_BG)
    run = cell.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) != str(PU._BRAND), "表头不该是红底/红字"


def test_title_is_red_and_body_stays_white():
    """标题红字 + 一条通栏红线，正文区不铺色块。

    原来是整条深红横幅压顶，一眼看过去最抢眼的是那块红。
    """
    pres = Presentation(PU.build_customer_status_pptx([_machine(0)]))
    slide = pres.slides[0]
    title = next(sh for sh in slide.shapes
                 if sh.has_text_frame and "客户面状态总览" in sh.text_frame.text)
    assert str(title.text_frame.paragraphs[0].runs[0].font.color.rgb) == str(PU._BRAND)
    # 红色**填充块**只有标题下那条线，且必须细——粗了就退化成横幅。
    # （文本框不带填充，红色在字上，所以这里数到的只会是色带。）
    red_bars = [sh for sh in slide.shapes if _solid_rgb(sh) == str(PU._BRAND)]
    assert len(red_bars) == 1, "红色块只应有标题下那一条线"
    assert red_bars[0].height / _EMU_IN <= 0.05


def test_every_page_carries_a_footer_inside_the_slide():
    """页脚每页都在，且不越过页底。

    导出的表经常被截图贴进别的材料，一张落单的图没有页码就找不回出处；
    而页脚一旦压出页底，PowerPoint 里是看不出报错的。
    """
    pres = Presentation(PU.build_customer_status_pptx([_machine(i, 3) for i in range(40)]))
    assert len(pres.slides) > 1
    for n, slide in enumerate(pres.slides, 1):
        texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        assert any(f"{n} /" in t for t in texts), f"第 {n} 页没有页码"
        assert any(PU._FOOTER_MARK in t for t in texts), f"第 {n} 页没有密级标识"
        for sh in slide.shapes:
            assert (sh.top + sh.height) / _EMU_IN <= _SLIDE_H_IN, \
                f"第 {n} 页有元素压出页底"


def test_title_block_does_not_overlap_the_table():
    """标题/副标题/红线都排在表格上方，一个都不许压到表头。"""
    pres = Presentation(PU.build_customer_status_pptx([_machine(0)]))
    slide = pres.slides[0]
    table_top = next(sh for sh in slide.shapes if sh.has_table).top / _EMU_IN
    for sh in slide.shapes:
        if sh.has_table or sh.top / _EMU_IN > 5:
            continue
        assert (sh.top + sh.height) / _EMU_IN <= table_top, "页头元素压到表格上了"
