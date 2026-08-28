"""PPT 导出工具：模板化 16:9 表格输出。

依赖：python-pptx

版式对齐部门述职 PPT 模板（见 CLAUDE.md「导出与上传」）：
- 白底正文；**红只出现在两处**——标题文字与标题下那条通栏横线
- 表头浅蓝灰底 + 黑字 + 居中加粗，分栏线用白色
- 数据行边框 + 左对齐 + 自动换行；斑马纹压到几乎看不见（13 列以上的表要靠它跟行）
- 状态点灯是整格底色（红黄绿灰），是一页里唯一的饱和色，所以一眼能扫到
- 每页固定页脚：口号 + 密级 + 页码
- 支持父子分组表头：通过将相邻列的父表头合并并设置同样的填充色实现
"""
import io
import json
import math
from datetime import datetime
from typing import Iterable, List, Optional, Sequence


def checklist_to_text(val: str) -> str:
    """将清单字段值（JSON 或旧纯文本）转为可读字符串，用于 PPT/导出。
    格式：每行前缀 ✓（已完成）或 ·（未完成）。
    """
    if not val:
        return ""
    try:
        items = json.loads(val)
        if isinstance(items, list):
            lines = []
            for item in items:
                text = str(item.get("text", "")).strip()
                if text:
                    lines.append(("✓ " if item.get("done") else "· ") + text)
            return "\n".join(lines)
    except (ValueError, TypeError, AttributeError):
        pass
    return val  # 旧纯文本，原样返回


def issues_to_text(machine, kind: str) -> str:
    """把机台的 customer_issues 条目转成 PPT 单元格文本。

    前缀：✓ 已闭环 / ⏸ 挂起 / · 进行中；挂起单独标出来，否则它在 PPT 里
    和未开始的看不出区别，评审时容易被当成"没人管"。
    已闭环的排在最后，重要的先入眼。
    kind="issue" 时把 demand（需求）一并带上（页面上两者同栏），加 [需求] 前缀区分。
    """
    kinds = {"issue", "demand"} if kind == "issue" else {kind}
    rows = [i for i in (getattr(machine, "issues", None) or []) if i.kind in kinds]
    if not rows:
        return ""
    rank = {"OPEN": 0, "挂起": 1, "CLOSED": 2}
    rows.sort(key=lambda i: (rank.get(i.status, 9), i.sort_order or 0, i.id))
    mark = {"CLOSED": "✓ ", "挂起": "⏸ "}
    lines = []
    for i in rows:
        text = (i.description or "").strip()
        if text:
            prefix = "[需求] " if i.kind == "demand" else ""
            lines.append(mark.get(i.status, "· ") + prefix + text)
    return "\n".join(lines)

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


_SLIDE_W = Emu(12192000)  # 16:9 默认 13.333"
_SLIDE_H = Emu(6858000)

# ── 配色：对齐部门述职 PPT 模板 ────────────────────────────────────────────
# 模板的取色逻辑是「红只用来定位、蓝只用来分层、饱和色只用来点灯」：
#   红   —— 只出现在标题文字与标题下那条横线上，一页里就这两处
#   浅蓝 —— 表头，把表头和正文分层
#   红黄绿 —— 只给状态格上底色，是整页唯一的饱和色，所以一眼就能扫到
# 原来的做法是整条深红横幅压顶 + 红底白字表头 + 浅红斑马，红铺满了半页，
# 结果最抢眼的是那块红，而看的人要找的是表里的数。
_BRAND = RGBColor(0xC7, 0x00, 0x0B)        # 华为红：标题文字 + 标题下横线
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_TEXT = RGBColor(0x26, 0x26, 0x26)
_MUTED = RGBColor(0x80, 0x80, 0x80)        # 副标题 / 页脚 / 页码
_HEADER_BG = RGBColor(0xD9, 0xE2, 0xF3)    # 表头：浅蓝灰
_HEADER_TEXT = RGBColor(0x1F, 0x24, 0x2E)  # 表头黑字（不是白字，底色浅）
_SECTION_BG = RGBColor(0xD9, 0xD9, 0xD9)   # 分组行 / 合计行：中灰
_BORDER = RGBColor(0xBF, 0xBF, 0xBF)       # 中性灰细边框
_RULE_H_IN = 0.028                         # 标题下横线厚度
# 斑马纹**保留但压到几乎看不见**，并且跟着表头挪到蓝灰一系。
# 模板里的表都在 6 列以内，白底就够认行；我们的产品需求表有 13~14 列，
# 全白的话眼睛横着扫一行会串到上下行去。压到这个亮度既不破坏版面、又够跟行。
_ZEBRA = RGBColor(0xF4, 0xF7, 0xFC)

# 中文 / 西文字体（华为优先 HarmonyOS Sans，回退微软雅黑）
_FONT_LATIN = "HarmonyOS Sans SC"
_FONT_EA = "微软雅黑"

# 进展状态点灯：**给格子上底色，不是给字上色**。
# 模板里的「状态点灯」「闭环状态」都是整格红黄绿，一屏扫过去哪几格是绿的一目了然；
# 只染字色的话，6 个进展列全是同一个字号的小字，得逐格去读才分得出来。
# 六档里只有四档给底色，「未开始 / 不涉及」故意留白底：
# 它们表达的是"这里没有进展"，上了底色反而和真有状态的格子一样抢眼。
_STATUS_FILLS = {
    "已完成": RGBColor(0x92, 0xD0, 0x50),   # 绿
    "进行中": RGBColor(0xFF, 0xD9, 0x66),   # 黄
    "已延期": RGBColor(0xFF, 0x99, 0x99),   # 红
    # 「已变更」＝这条需求本轮不做了。导出**不剔这些行**（那是交付记录），
    # 但灰底把它和还在推进的行区分开，看的人不会把它算进进度里。
    "已变更": RGBColor(0xD9, 0xD9, 0xD9),   # 灰
}
_STATUS_TEXT = {
    "未开始": RGBColor(0x90, 0x93, 0x99),
    "不涉及": RGBColor(0xB0, 0xB3, 0xB8),
}
# 点灯底色都是浅色，字一律用正文黑；白字在黄底上等于没有。
_STATUS_ON_FILL_TEXT = RGBColor(0x1F, 0x24, 0x2E)

# 页脚：模板每页固定带的三件套（口号 / 密级 / 页码）。
# 不想要就把这两个常量置空，页脚只剩分隔线与页码。
_FOOTER_BRAND = "HILIGHT"
_FOOTER_SLOGAN = "求真务实 · 攻坚克难 · 开放协同 · 整体最优"
_FOOTER_MARK = "Restricted Distribution"


# ─── 版面常量：一张 16:9 幻灯片里表格能占多大 ───────────────────────────────
# 幻灯片 13.333" × 7.5"；横幅 0.8"，页脚留 0.35"。
_MARGIN_IN = 0.4
_TABLE_W_IN = 12.53          # 13.333 - 2×0.4
_TABLE_TOP_IN = 1.0
_TABLE_BOTTOM_IN = 7.02      # 再往下就压到页脚了
_TABLE_H_IN = _TABLE_BOTTOM_IN - _TABLE_TOP_IN
_FOOTER_TOP_IN = 7.10        # 页脚分隔线，表格底边再往下 0.08"

_CELL_PAD_LR_PT = 5          # 与 _set_cell_margins 一致
_CELL_PAD_TB_PT = 2
_LINE_FACTOR = 1.25          # 行距系数（中文字面高 + 行间）
_MIN_ROW_H_IN = 0.26

# 单元格最多渲染几行，超出截断。**一格不能吃掉一整页**：客户面的「关键事务」
# 一台机器挂十几条，不截断的话那一行能顶满整张幻灯片，后面全被挤到页外。
_DEFAULT_MAX_LINES = 6


def _text_em_width(s: str) -> float:
    """一行文字的宽度，单位 em（1 em ＝ 字号）。

    中日韩全角字按 1.0 em、其余按 0.52 em 估。这是**估算**，宁可估宽一点：
    估窄了会算出"装得下"，导出的表就又超出页面了，而那正是要修的问题。
    """
    w = 0.0
    for ch in s:
        w += 1.0 if ord(ch) > 0x2E7F else 0.52
    return w


def _wrapped_lines(text: str, col_w_in: float, font_pt: int) -> int:
    """这段文字在这么宽的列里会占几行（含手动换行）。"""
    avail_in = max(col_w_in - 2 * _CELL_PAD_LR_PT / 72.0, 0.2)
    per_line_em = avail_in * 72.0 / font_pt
    lines = 0
    for seg in (text or "").split("\n"):
        lines += max(1, math.ceil(_text_em_width(seg) / per_line_em)) if seg else 1
    return max(lines, 1)


def _clip_text(text: str, col_w_in: float, font_pt: int, max_lines: int) -> str:
    """按渲染后的行数截断，并如实写明还剩多少行。

    直接截字符会把最后一条截成半句；按行截并标出"另 N 条"，看的人知道
    PPT 里不是全部，要全量去系统里看——比悄悄少几行强。
    """
    segs = (text or "").split("\n")
    if len(segs) <= 1 and _wrapped_lines(text, col_w_in, font_pt) <= max_lines:
        return text
    out, used = [], 0
    for i, seg in enumerate(segs):
        need = _wrapped_lines(seg, col_w_in, font_pt)
        if used + need > max_lines - (1 if i < len(segs) - 1 else 0):
            rest = len(segs) - i
            out.append(f"…（另 {rest} 条，见系统）")
            return "\n".join(out)
        out.append(seg)
        used += need
    return "\n".join(out)


def _row_height_in(values: Sequence[str], widths: Sequence[float], font_pt: int) -> float:
    lines = 1
    for val, w in zip(values, widths):
        lines = max(lines, _wrapped_lines("" if val is None else str(val), w, font_pt))
    h = lines * font_pt * _LINE_FACTOR / 72.0 + 2 * _CELL_PAD_TB_PT / 72.0
    return max(h, _MIN_ROW_H_IN)


def _norm_widths(widths: Optional[Sequence[float]], n_cols: int) -> List[float]:
    """把列宽等比缩放到正好铺满表格宽度。

    原来是手写一串英寸数直接用：加一列、改个标题就再也对不上，
    表格要么越出页面右边、要么右边空一大条。归一化之后，列宽只表达**比例**。
    """
    if not widths or len(widths) != n_cols:
        widths = [1.0] * n_cols
    total = sum(widths) or 1.0
    return [w * _TABLE_W_IN / total for w in widths]


def _fit_font_size(n_cols: int) -> int:
    """字号只跟**列数**走，不跟行数走。

    行数由分页解决——原来行一多就把字号压到 8pt，结果是"既看不清、又还是超出页面"。
    """
    if n_cols <= 8:
        return 11
    if n_cols <= 12:
        return 10
    if n_cols <= 16:
        return 9
    return 8


def _paginate(rows: Sequence[Sequence[str]], widths: Sequence[float],
              font_pt: int, header_h_in: float) -> List[List[Sequence[str]]]:
    """按**估算行高**把数据行切成若干页，每页都装得下。

    不是按固定条数切：客户面一行可能是 6 行文字、也可能是 1 行，
    按条数切的结果是有的页空半张、有的页照样溢出。
    """
    budget = _TABLE_H_IN - header_h_in
    pages: List[List[Sequence[str]]] = []
    cur: List[Sequence[str]] = []
    used = 0.0
    for row in rows:
        h = _row_height_in(row, widths, font_pt)
        if cur and used + h > budget:
            pages.append(cur)
            cur, used = [], 0.0
        cur.append(row)
        used += h
    if cur:
        pages.append(cur)
    return pages or [[]]


def _new_pres() -> Presentation:
    pres = Presentation()
    pres.slide_width = _SLIDE_W
    pres.slide_height = _SLIDE_H
    return pres


def _add_bar(slide, top_in: float, height_in: float, color: RGBColor):
    """一条通栏色带（标题下的红线、页脚上的灰线）。用矩形而不是直线：
    直线在 PowerPoint 里会带上主题的线端/阴影，通栏时两头会翘。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top_in),
                                 _SLIDE_W, Inches(height_in))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.shadow.inherit = False
    return bar


def _add_header(slide, title: str, subtitle: str):
    """页头：红色标题 + 灰色副标题 + 一条通栏红线，正文区全白。

    标题与副标题**分两行**排，不并排：副标题是「导出时间 · 共 N 条 · 第 X/Y 页」，
    一长起来就会撞上左边的标题，而标题是不换行的（一折行就压到表格上）。
    """
    title_box = slide.shapes.add_textbox(Inches(_MARGIN_IN), Inches(0.14),
                                         Inches(_TABLE_W_IN), Inches(0.42))
    tf = title_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # 不换行：标题一折行就压到表格上；宁可长标题右边留白，也不要两行标题顶掉表头
    tf.word_wrap = False
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _apply_run_font(p.runs[0], 22, True, _BRAND)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(_MARGIN_IN), Inches(0.53),
                                           Inches(_TABLE_W_IN), Inches(0.26))
        tf2 = sub_box.text_frame
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        tf2.word_wrap = False
        tf2.text = subtitle
        _apply_run_font(tf2.paragraphs[0].runs[0], 10, False, _MUTED)

    _add_bar(slide, 0.82, _RULE_H_IN, _BRAND)


def _rgb_to_hex(color: RGBColor) -> str:
    return str(color)  # python-pptx RGBColor 的 __str__ 返回 6 位大写十六进制


def _apply_run_font(run, size: int, bold: bool, color: RGBColor):
    """统一设置 run 的字号/粗细/颜色，并补齐东亚字体（中文不走默认衬线）。"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = _FONT_LATIN
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", _FONT_EA)


def _set_cell_margins(cell, lr=5, tb=2):
    cell.margin_left = Pt(lr)
    cell.margin_right = Pt(lr)
    cell.margin_top = Pt(tb)
    cell.margin_bottom = Pt(tb)


def _clear_table_style(table):
    """清掉 python-pptx 默认套用的主题表样式（带蓝色条纹），
    换成"无样式"，让我们手动设置的填充/边框完全生效。"""
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn("a:tblPr"))
    for attr in ("firstRow", "lastRow", "firstCol", "lastCol", "bandRow", "bandCol"):
        tblPr.set(attr, "0")
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is None:
        style_id = etree.SubElement(tblPr, qn("a:tableStyleId"))
    # "No Style, No Grid"
    style_id.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _set_cell_border(cell, color: RGBColor = _BORDER):
    """给单元格四边加细边框。python-pptx 没有现成 API，用 lxml 操作。"""
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        # 先删除已有，避免重复
        existing = tc_pr.findall(qn(tag))
        for e in existing:
            tc_pr.remove(e)
    hex_val = _rgb_to_hex(color)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = etree.SubElement(tc_pr, qn(tag))
        ln.set("w", "6350")  # 0.5pt
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        solid = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(solid, qn("a:srgbClr"))
        srgb.set("val", hex_val)


def _style_header_cell(cell, text: str, font_size: int):
    """表头：浅蓝灰底 + 黑字。**不是红底白字**——一页里的红要留给标题，
    表头一红，整张表的重心就压在最上面一行，正文反而看不见。
    分隔线用白色：在浅蓝底上白线是"分栏"，灰线会糊成一片。"""
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = _HEADER_BG
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_cell_margins(cell)
    for para in cell.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for r in para.runs:
            _apply_run_font(r, font_size + 1, True, _HEADER_TEXT)
    _set_cell_border(cell, _WHITE)


_TOTAL_LABELS = {"合计", "总计", "小计"}
_TOTAL_BG = _SECTION_BG      # 合计行＝中灰，与模板里的分组行同一档


_LINK_COLOR = RGBColor(0x15, 0x65, 0xC0)


def _style_link_cell(cell, url: str, font_size: int, zebra: bool, label: str = "查看"):
    """URL 列渲染成可点的短链接。

    整条 URL 摊在格子里既把列撑宽、又没人会照着念——PPT 里能点开才有用。
    """
    cell.text = ""
    cell.fill.solid()
    cell.fill.fore_color.rgb = _ZEBRA if zebra else _WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_cell_margins(cell)
    para = cell.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    _apply_run_font(run, font_size, False, _LINK_COLOR)
    run.font.underline = True
    run.hyperlink.address = url
    _set_cell_border(cell)


def _style_data_cell(cell, value, font_size: int, zebra: bool, center: bool = False,
                     total: bool = False):
    text = "" if value is None else str(value)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = _TOTAL_BG if total else (_ZEBRA if zebra else _WHITE)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_cell_margins(cell)

    # 进展状态点灯：有底色的四档给格子上底色，另两档只压字色。
    # 合计行不点灯——它的灰底是"这一行是汇总"，被状态色盖掉就认不出来了。
    stripped = text.strip()
    fill = None if total else _STATUS_FILLS.get(stripped)
    if fill is not None:
        cell.fill.fore_color.rgb = fill
    muted = _STATUS_TEXT.get(stripped)
    is_status = fill is not None or muted is not None
    color = _STATUS_ON_FILL_TEXT if fill is not None else (muted or _TEXT)

    for para in cell.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER if (center or is_status) else PP_ALIGN.LEFT
        for r in para.runs:
            _apply_run_font(r, font_size, (fill is not None) or total, color)
    _set_cell_border(cell)


def _merge_header_groups(table, header_row: int, groups: Sequence[tuple]):
    """合并某一表头行的相邻列：groups = [(start_col, end_col, label), ...]"""
    for start, end, label in groups:
        if end > start:
            table.cell(header_row, start).merge(table.cell(header_row, end))
        cell = table.cell(header_row, start)
        cell.text = label


def _add_grouped_table(
    slide,
    parent_headers: Sequence[Optional[tuple]],
    leaf_headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    widths: Sequence[float],
    font_size: int,
    center_cols: Optional[set] = None,
    link_cols: Optional[set] = None,
):
    """在一张 slide 上画一张「父表头 + 子表头」的表。**只画一页**，分页由调用方做。

    parent_headers: 长度与 leaf_headers 相同。每个元素是 (group_id, label) 或 None。
        相邻同 group_id 会被合并；None 表示该列没有父分组，与子表头垂直合并。
    widths: 已经归一化过的列宽（英寸），加起来正好是表格宽度。
    """
    center_cols = center_cols or set()
    link_cols = link_cols or set()
    n_cols = len(leaf_headers)
    n_rows = 2 + len(rows)

    header_h = _row_height_in(leaf_headers, widths, font_size + 1)
    body_h = sum(_row_height_in(r, widths, font_size) for r in rows)
    total_h = min(2 * header_h + body_h, _TABLE_H_IN)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(_MARGIN_IN), Inches(_TABLE_TOP_IN),
        Inches(_TABLE_W_IN), Inches(total_h))
    table = table_shape.table
    _clear_table_style(table)

    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    # 行高显式写死成估算值：不写的话 add_table 会把总高均分给每一行，
    # 短行被撑得老高，长行还是溢出，整张表看着松松垮垮
    table.rows[0].height = Inches(header_h)
    table.rows[1].height = Inches(header_h)
    for i, row in enumerate(rows, start=2):
        table.rows[i].height = Inches(_row_height_in(row, widths, font_size))

    # ===== 父表头行 =====
    for j in range(n_cols):
        _style_header_cell(table.cell(0, j), "", font_size)

    groups: List[tuple] = []
    j = 0
    while j < n_cols:
        ph = parent_headers[j]
        if ph is None:
            j += 1
            continue
        gid, label = ph
        end = j
        while end + 1 < n_cols and parent_headers[end + 1] is not None and parent_headers[end + 1][0] == gid:
            end += 1
        groups.append((j, end, label, gid))
        j = end + 1

    _merge_header_groups(table, 0, [(s_, e_, lab) for s_, e_, lab, _ in groups])
    for s_, e_, lab, _ in groups:
        _style_header_cell(table.cell(0, s_), lab, font_size)

    # ===== 子表头行 =====
    for j, h in enumerate(leaf_headers):
        _style_header_cell(table.cell(1, j), h, font_size)

    # 没有父分组的列：父行与子行垂直合并
    for j in range(n_cols):
        if parent_headers[j] is None:
            table.cell(0, j).merge(table.cell(1, j))
            _style_header_cell(table.cell(0, j), leaf_headers[j], font_size)

    # ===== 数据行 =====
    for i, row in enumerate(rows, start=2):
        zebra = (i - 2) % 2 == 1
        # 合计行加粗 + 浅底：混在斑马纹里根本认不出来哪一行是合计
        total = bool(row) and str(row[0]).strip() in _TOTAL_LABELS
        for j, val in enumerate(row):
            text = "" if val is None else str(val)
            if j in link_cols and text.startswith(("http://", "https://")):
                _style_link_cell(table.cell(i, j), text, font_size, zebra)
            else:
                _style_data_cell(table.cell(i, j), val, font_size, zebra,
                                 center=(j in center_cols), total=total)


def _footer_text(slide, left_in: float, width_in: float, align, text: str,
                 size: int = 8, color: RGBColor = _MUTED, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(_FOOTER_TOP_IN + 0.05),
                                   Inches(width_in), Inches(0.26))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    # text 可能是空串（调用方随后自己 add_run 拼多段），空段落里没有 run
    if para.runs:
        _apply_run_font(para.runs[0], size, bold, color)
    return para


def _add_footer(slide, idx: int, total: int):
    """页脚：细灰分隔线 + 左口号 + 中密级 + 右页码。

    页码**每页都写**，不只在多页时写：导出的表经常被截图贴进别的材料，
    一张落单的图没有页码就找不回它是第几页、更找不回是哪一次导出的。
    """
    _add_bar(slide, _FOOTER_TOP_IN, 0.012, _SECTION_BG)

    if _FOOTER_BRAND or _FOOTER_SLOGAN:
        para = _footer_text(slide, _MARGIN_IN, 7.0, PP_ALIGN.LEFT, "")
        if _FOOTER_BRAND:
            run = para.add_run()
            run.text = _FOOTER_BRAND + "  "
            _apply_run_font(run, 9, True, _BRAND)
            run.font.italic = True
        if _FOOTER_SLOGAN:
            run = para.add_run()
            run.text = _FOOTER_SLOGAN
            _apply_run_font(run, 8, False, _MUTED)

    if _FOOTER_MARK:
        _footer_text(slide, _MARGIN_IN, _TABLE_W_IN, PP_ALIGN.CENTER, _FOOTER_MARK)

    _footer_text(slide, 10.6, 2.33, PP_ALIGN.RIGHT,
                 f"{idx} / {total}" if total > 1 else str(idx), size=9)


def _add_empty_slide(pres, title: str, subtitle: str, note: str):
    """一条数据都没有时给一张明确写着"暂无数据"的页。

    原来会生成一张只有表头的空表，看着像导出坏了——而它其实是对的。
    """
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    _add_header(slide, title, subtitle)
    box = slide.shapes.add_textbox(Inches(_MARGIN_IN), Inches(3.0), Inches(_TABLE_W_IN), Inches(0.6))
    tf = box.text_frame
    tf.text = note
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _apply_run_font(p.runs[0], 16, False, _MUTED)
    _add_footer(slide, 1, 1)
    return slide


def add_table_slides(
    pres,
    title: str,
    subtitle: str,
    leaf_headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    parent_headers: Optional[Sequence[Optional[tuple]]] = None,
    col_ratios: Optional[Sequence[float]] = None,
    center_cols: Optional[set] = None,
    clip_cols: Optional[dict] = None,
    link_cols: Optional[set] = None,
    empty_note: str = "暂无数据",
):
    """把一批数据行铺成**若干张**幻灯片，每张都装得下。

    这是所有 PPT 表格的唯一入口。原来是"一张 slide 一张表，有多少行画多少行"，
    行一多表格就顺着页面往下长、直接长到幻灯片外面去——导出的 PPT 拿到手要么
    整段看不见、要么得手工拆页，这正是"格式没办法直接用"的主因。

    col_ratios: 列宽**比例**（不是英寸），内部归一化到正好铺满页宽。
    clip_cols:  {列下标: 最多行数}，长文本列超出就截断并标"…（另 N 条）"。
                一格能吃掉一整页，不截断的话分页也救不回来。
    link_cols:  这些列里的 http(s) 值渲染成可点的「查看」短链接。
    """
    n_cols = len(leaf_headers)
    parent_headers = list(parent_headers or [None] * n_cols)
    widths = _norm_widths(col_ratios, n_cols)
    font_size = _fit_font_size(n_cols)

    clip_cols = clip_cols or {}
    link_cols = link_cols or set()
    prepared: List[List[str]] = []
    for row in rows:
        vals = ["" if v is None else str(v) for v in row]
        for j, max_lines in clip_cols.items():
            if j < len(vals):
                vals[j] = _clip_text(vals[j], widths[j], font_size, max_lines)
        prepared.append(vals)

    if not prepared:
        _add_empty_slide(pres, title, subtitle, empty_note)
        return

    header_h = _row_height_in(leaf_headers, widths, font_size + 1)
    # 链接列渲染出来只有「查看」两个字，估行高时不能拿整条 URL 去算：
    # 一条长链接会把整行的估算高度顶上去，白白少排好几行
    est_rows = [[("查看" if (j in link_cols and v.startswith("http")) else v)
                 for j, v in enumerate(r)] for r in prepared]
    sizes = [len(pg) for pg in _paginate(est_rows, widths, font_size, header_h * 2)]
    pages: List[List[str]] = []
    cursor = 0
    for n in sizes:
        pages.append(prepared[cursor:cursor + n])
        cursor += n
    for idx, page_rows in enumerate(pages, start=1):
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        page_sub = subtitle
        if len(pages) > 1:
            page_sub += f"   ·   第 {idx}/{len(pages)} 页"
        _add_header(slide, title, page_sub)
        _add_grouped_table(slide, parent_headers, leaf_headers, page_rows,
                           widths, font_size, center_cols, link_cols)
        _add_footer(slide, idx, len(pages))


def build_customer_status_pptx(rows: Iterable) -> io.BytesIO:
    """传入 CustomerStatus ORM 列表，返回 BytesIO。"""
    leaf_headers = [
        "机台编号", "客户", "型号", "当前阶段", "现场版本", "关注度",
        "当前进展", "现场关键事务", "软件类风险和问题", "问题单",
    ]
    # 列宽是**比例**，内部归一化到页宽：手写英寸数的话，改个列就再也铺不满
    col_ratios = [0.85, 1.0, 0.9, 0.95, 0.95, 0.7, 2.0, 2.0, 2.0, 1.15]

    data: List[List[str]] = []
    for r in rows:
        data.append([
            r.machine_id or "",
            r.battlefield or "",
            r.model or "",
            r.current_stage or "",
            r.field_version or "",
            ("★" * (r.attention_level or 0)) or "-",
            r.customer_status or "",
            issues_to_text(r, "task"),
            issues_to_text(r, "issue"),
            getattr(r, "issue_url", "") or "—",
        ])

    pres = _new_pres()
    add_table_slides(
        pres,
        "客户面状态总览",
        f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}   ·   共 {len(data)} 台机器",
        leaf_headers, data,
        col_ratios=col_ratios,
        center_cols={0, 2, 3, 4, 5, 9},
        link_cols={9},
        # 三个长文本列各限 6 行：一台机器挂十几条事务的话，
        # 不截断那一行能顶满整张幻灯片，后面的机器全被挤到页外
        clip_cols={6: _DEFAULT_MAX_LINES, 7: _DEFAULT_MAX_LINES, 8: _DEFAULT_MAX_LINES},
        empty_note="没有可导出的机台",
    )

    buf = io.BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf


def _iteration_title(iteration) -> str:
    title = f"{iteration.year}年{iteration.month}月迭代"
    if iteration.name:
        title += f"  ·  {iteration.name}"
    return title


def _add_domain_slide(pres, iteration, requirements: Iterable):
    """领域需求：基础列 + 「交付进展跟踪」6 子列 + 备注，超过一页自动续页。"""
    leaf_headers = [
        "序号", "需求编号", "需求标题", "责任人", "PL组", "优先级", "计划版本",
        "需求串讲", "反串讲", "STC设计", "编码", "BBIT", "转测澄清",
        "备注",
    ]
    parent_headers: List[Optional[tuple]] = [None] * 7 + [("progress", "交付进展跟踪")] * 6 + [None]
    col_ratios = [0.45, 1.0, 2.0, 0.75, 0.7, 0.55, 0.85] + [0.82] * 6 + [1.4]

    data: List[List[str]] = []
    for r in requirements:
        data.append([
            str(r.seq or 0),
            r.req_no or "",
            r.title or "",
            r.owner or "",
            getattr(r, "owner_group", "") or "",
            r.priority or "",
            r.planned_version or "",
            r.progress_walkthrough or "",
            r.progress_reverse or "",
            r.progress_stc or "",
            r.progress_coding or "",
            r.progress_bbit or "",
            r.progress_clarify or "",
            getattr(r, "remark", "") or "",
        ])

    subtitle = f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}   ·   共 {len(data)} 条领域需求"
    if iteration.owner:
        subtitle += f"   ·   负责人：{iteration.owner}"
    add_table_slides(
        pres, _iteration_title(iteration) + "  ·  领域需求", subtitle,
        leaf_headers, data,
        parent_headers=parent_headers, col_ratios=col_ratios,
        center_cols={0, 5, 7, 8, 9, 10, 11, 12},
        clip_cols={2: 4, 13: 4},
        empty_note="本迭代还没有领域需求",
    )


def _add_product_slides(pres, iteration, product_reqs: Iterable):
    """产品需求：单页字段太多，拆成「基础信息」与「交付进展跟踪」两组，各自按需续页。"""
    rows_basic: List[List[str]] = []
    rows_progress: List[List[str]] = []
    for r in product_reqs:
        rows_basic.append([
            str(r.seq or 0),
            r.req_no or "",
            r.title or "",
            r.planned_version or "",
            r.priority or "",
            r.feature or "",
            r.feature_fo or "",
            r.feature_se or "",
            r.feature_tfo or "",
            r.code_areas or "",
            r.key_risks or "",
        ])
        rows_progress.append([
            str(r.seq or 0),
            r.req_no or "",
            r.title or "",
            r.progress_walkthrough or "",
            r.progress_reverse or "",
            r.progress_domain or "",
            r.progress_coding or "",
            r.progress_joint_debug or "",
            r.progress_clarify or "",
            r.progress_test_result or "",
            r.estimated_loc or "",
            r.actual_loc or "",
            r.actual_effort or "",
        ])

    base_title = _iteration_title(iteration)
    subtitle = f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}   ·   共 {len(rows_basic)} 条产品需求"

    add_table_slides(
        pres, base_title + "  ·  产品需求（基础信息）", subtitle,
        ["序号", "需求编号", "需求标题", "计划版本", "优先级", "所属特性",
         "特性FO", "特性SE", "特性TFO", "涉及代码领域", "关键风险"],
        rows_basic,
        col_ratios=[0.4, 1.0, 1.9, 0.9, 0.55, 1.0, 0.7, 0.7, 0.7, 1.6, 2.5],
        center_cols={0, 4, 6, 7, 8},
        clip_cols={2: 4, 9: 5, 10: 5},
        empty_note="本迭代还没有产品需求",
    )

    add_table_slides(
        pres, base_title + "  ·  产品需求（交付进展跟踪）", subtitle,
        ["序号", "需求编号", "需求标题",
         "需求串讲", "反串讲", "领域串讲", "编码", "联调验证", "转测澄清", "测试结论",
         "预估代码量", "实际代码量", "实际工作量"],
        rows_progress,
        parent_headers=[None, None, None] + [("progress", "交付进展跟踪")] * 7 + [None] * 3,
        col_ratios=[0.4, 1.0, 1.5, 0.8, 0.7, 0.85, 0.6, 0.85, 0.85, 0.85, 0.95, 0.95, 0.95],
        center_cols={0, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12},
        clip_cols={2: 4},
        empty_note="本迭代还没有产品需求",
    )


def build_iteration_pptx(iteration, requirements: Iterable, product_reqs: Optional[Iterable] = None) -> io.BytesIO:
    """传入 AnnualIteration + 领域/产品 需求列表，返回 BytesIO。

    - 始终输出领域需求 slide（保持原有行为）
    - 若 product_reqs 非空，再追加 2 张产品需求 slide
    """
    pres = _new_pres()
    _add_domain_slide(pres, iteration, requirements)

    if product_reqs:
        product_list = list(product_reqs)
        if product_list:
            _add_product_slides(pres, iteration, product_list)

    buf = io.BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf


# 每个数值列至少这么宽，否则表头「2026-03」会被压成两行、客户名直接看不清
_MIN_MATRIX_COL_IN = 0.62


def add_matrix_slides(pres, title: str, subtitle: str, label_header: str,
                      value_headers: Sequence[str], rows: Sequence[Sequence],
                      empty_note: str = "暂无数据"):
    """宽矩阵表（小组 × 月份 / 客户 / 特性）：**列太多就按列切页**，每页都带上标签列。

    原来是把所有列平分页宽：二十几个客户时每列只剩 0.4"，表头挤成竖排，
    数字缩成一团——"导出来没法直接用"的另一半原因。
    """
    n_val = len(value_headers)
    per = max(1, int((_TABLE_W_IN - 1.9) // _MIN_MATRIX_COL_IN))
    if not n_val:
        _add_empty_slide(pres, title, subtitle, empty_note)
        return
    for start in range(0, n_val, per):
        chunk = list(value_headers[start:start + per])
        heads = [label_header] + chunk
        sub_rows = [[r[0]] + list(r[1 + start:1 + start + len(chunk)]) for r in rows]
        t = title
        if n_val > per:
            t += f"（{start + 1}-{start + len(chunk)} 列，共 {n_val} 列）"
        add_table_slides(
            pres, t, subtitle, heads, sub_rows,
            col_ratios=[2.0] + [1.0] * len(chunk),
            center_cols=set(range(1, len(heads))),
            empty_note=empty_note,
        )
