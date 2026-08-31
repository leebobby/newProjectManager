"""问题单报表接口：读取服务器侧 Excel，解析后返回给前端。

权限：
- GET  /api/issues/data              —— 所有登录用户（加载最新单张报表）
- GET  /api/issues/trend             —— 所有登录用户（扫描全目录按天聚合趋势）
- GET  /api/issues/snapshot-flow     —— 所有登录用户（每日新增/解决，相邻快照差分）
- GET  /api/issues/flow-detail       —— 所有登录用户（某天新增/解决的明细）
- GET  /api/issues/ungrouped         —— 所有登录用户（最新快照里归不到小组的责任人）
- GET  /api/issues/run-script/status —— 所有登录用户（查询脚本是否正在运行）
- POST /api/issues/run-script        —— 仅管理员（执行外部刷新脚本）
- GET  /api/issues/export.pptx       —— 所有登录用户（导出 PPT）

配置项（via PUT /api/config）：
  issue_report_path      —— Excel 目录或文件路径
  issue_script_path      —— 刷新脚本路径（.py / .bat / .exe）
  issue_script_timeout   —— 采集脚本超时秒数（默认 600）
  issue_snapshot_time    —— 每日自动采集时刻 HH:MM（默认 07:30）
  issue_snapshot_enabled —— 是否启用每日自动采集（默认 true）
"""
import io
import json
import pathlib
import re
import subprocess
import sys
import threading
from datetime import datetime
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session

import models
from auth import get_current_user, require_admin
from database import SessionLocal, get_db
from op_log import log_op
from routers.config import _load as _load_config
from timeutil import fmt_local, iso_local

router = APIRouter(prefix="/api/issues", tags=["issues"])

# 全局锁：防止脚本并发执行
_script_lock: threading.Lock = threading.Lock()
_script_started_at: Optional[datetime] = None

_RAW_COLS = [
    "version",          # A 版本信息
    "issue_id",         # B 缺陷业务编号
    "title",            # C 标题
    "owner",            # D 当前责任人
    "group",            # E 当前责任人所属小组
    "progress",         # F 进展
    "severity",         # G 严重程度
    "severity_di",      # H 严重程度DI值
    "root_cause",       # I 根因
    "solution",         # J 解决措施
    "progress_record",  # K 进展记录
    "estimated_close",  # L 预计闭环时间
    "priority",         # M 优先级
    "is_sdts",          # N 是否SDTS
    "year",             # O 年份
    "month",            # P 月份
    "date",             # Q 日期
    "year_month",       # R 年月（钻取按月度过滤的关键字段）
    "category",         # S 标题分类（钻取按客户/分类过滤的关键字段）
    "customer",         # T 客户面（API 快照聚合的关键维度；由后端从标题匹配客户主数据得到）
    "department",       # U 责任人部门（展示用：责任人所在的直属部门）
    "feature",          # V 特性
    "subsystem",        # W 子系统
    "module",           # X 模块
    "dept_path",        # Y 责任人部门全路径（各级部门拼接，仅用于部门过滤匹配，不展示）
]

_DATE_PAT      = re.compile(r"_(\d{8})\.",           re.IGNORECASE)
_DATE_DIR_PAT  = re.compile(r"^\d{4}-\d{2}-\d{2}$")

COLORS = ["#4073ba", "#67C23A", "#E6A23C", "#F56C6C", "#909399", "#8E7AD8", "#26C9C3"]

# 责任人不在任何小组名单时的归属。**只有一个字面量**：Excel 交叉表、维度聚合、
# 前端提示都用它。两处各写一个（曾经是「未分组」/「未归组」）的表现是同一批人
# 在两张表里分成两档，加起来还对，看着都像对的。
UNGROUPED_GROUP = "未归组"


# ─── helpers ────────────────────────────────────────────────────────────────

def _cell(ws, row: int, col: int) -> str:
    v = ws.cell(row, col).value
    return str(v).strip() if v is not None else ""


def _count_by(rows: List[Dict], field: str) -> Dict[str, int]:
    result: Dict[str, int] = {}
    for r in rows:
        k = r.get(field, "")
        if k:
            result[k] = result.get(k, 0) + 1
    return result


def _file_sort_key(f: pathlib.Path):
    m = _DATE_PAT.search(f.name)
    if m:
        try:
            return (1, datetime.strptime(m.group(1), "%Y%m%d").timestamp())
        except ValueError:
            pass
    return (0, f.stat().st_mtime)


def _resolve_target(path_str: str) -> pathlib.Path:
    p = pathlib.Path(path_str)
    if not p.exists():
        raise HTTPException(404, f"路径不存在：{path_str}")
    if p.is_file():
        if p.suffix.lower() != ".xlsx":
            raise HTTPException(400, "指定文件不是 .xlsx 格式")
        return p
    if p.is_dir():
        candidates = sorted(p.glob("*.xlsx"), key=_file_sort_key, reverse=True)
        if not candidates:
            raise HTTPException(404, f"目录 {path_str} 中未找到 .xlsx 文件")
        return candidates[0]
    raise HTTPException(400, f"无法识别路径：{path_str}")


def _parse_cross_table(ws) -> Dict[str, Any]:
    max_col = ws.max_column
    columns = [_cell(ws, 1, c) for c in range(2, max_col + 1)]
    while columns and not columns[-1]:
        columns.pop()
    rows: List[Dict] = []
    for r in range(2, ws.max_row + 1):
        label = _cell(ws, r, 1)
        if not label:
            break
        row: Dict[str, Any] = {"label": label}
        for i, col in enumerate(columns):
            raw_val = ws.cell(r, i + 2).value
            if raw_val is None:
                row[col] = 0
            else:
                try:
                    row[col] = int(raw_val)
                except (ValueError, TypeError):
                    row[col] = str(raw_val)
        rows.append(row)
    return {"columns": columns, "rows": rows}


# 最新报表解析缓存：str(path) -> (mtime, 解析结果)。文件没变就不重复开 openpyxl，
# 让"打开问题单管理→看最新"这条最常走的路径秒回。进程内、按文件粒度。
_parse_cache: Dict[str, tuple] = {}


def _parse_excel_cached(path: str) -> Dict[str, Any]:
    """带 mtime 缓存的报表解析。文件未变直接命中缓存，否则重新解析并刷新缓存。"""
    p = pathlib.Path(path)
    try:
        mtime = p.stat().st_mtime
    except OSError:
        mtime = None
    key = str(p)
    hit = _parse_cache.get(key)
    if hit is not None and mtime is not None and hit[0] == mtime:
        return hit[1]
    result = _parse_excel(path)
    if mtime is not None:
        _parse_cache[key] = (mtime, result)
    return result


def _parse_excel(path: str) -> Dict[str, Any]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
    except FileNotFoundError:
        raise HTTPException(404, "报表文件不存在，请检查路径配置")
    except Exception as exc:
        raise HTTPException(500, f"读取报表文件失败: {exc}")

    raw: List[Dict] = []
    try:
        ws_raw = wb["原始数据"]
        for r in range(2, ws_raw.max_row + 1):
            row = {col: _cell(ws_raw, r, i + 1) for i, col in enumerate(_RAW_COLS)}
            if any(row.values()):
                raw.append(row)
    except KeyError:
        pass

    def _sheet(name):
        try:
            return _parse_cross_table(wb[name])
        except KeyError:
            return {"columns": [], "rows": []}

    try:
        mtime = datetime.fromtimestamp(
            pathlib.Path(path).stat().st_mtime
        ).strftime("%Y-%m-%d %H:%M:%S")
    except Exception:
        mtime = None

    return {
        "file_mtime": mtime,
        "raw": raw,
        "monthly_by_group":   _sheet("按小组月度统计"),
        "annual_by_group":    _sheet("按小组年度统计"),
        "by_customer":        _sheet("按客户统计"),
        "feature_by_group":   _sheet("特性×小组统计"),
        "feature_by_customer":_sheet("特性×客户统计"),
    }


def _parse_raw_from_wb(wb) -> List[Dict]:
    raw: List[Dict] = []
    try:
        ws = wb["原始数据"]
        for r in range(2, ws.max_row + 1):
            row = {col: _cell(ws, r, i + 1) for i, col in enumerate(_RAW_COLS)}
            if any(row.values()):
                raw.append(row)
    except KeyError:
        pass
    return raw


def _list_date_dirs(root: pathlib.Path) -> list:
    """Return date subdirs (YYYY-MM-DD) containing xlsx files, sorted newest-first."""
    result = [
        d for d in root.iterdir()
        if d.is_dir() and _DATE_DIR_PAT.match(d.name) and any(d.glob("*.xlsx"))
    ]
    return sorted(result, key=lambda d: d.name, reverse=True)


def _resolve_for_date(path_str: str, date: Optional[str] = None):
    """Resolve to a single xlsx path, supporting both flat dirs and date-subdir structures."""
    p = pathlib.Path(path_str)
    if not p.exists():
        raise HTTPException(404, f"路径不存在：{path_str}")
    if p.is_file():
        return p

    if p.is_dir():
        date_dirs = _list_date_dirs(p)
        if date_dirs:
            if date:
                target_dir = p / date
                if not target_dir.is_dir():
                    raise HTTPException(404, f"日期目录不存在：{date}")
                xlsxes = sorted(target_dir.glob("*.xlsx"), key=lambda f: f.name, reverse=True)
                if not xlsxes:
                    raise HTTPException(404, f"日期 {date} 目录中无 xlsx 文件")
            else:
                target_dir = date_dirs[0]
                xlsxes = sorted(target_dir.glob("*.xlsx"), key=lambda f: f.name, reverse=True)
            return xlsxes[0]

        # Flat directory fallback
        candidates = sorted(p.glob("*.xlsx"), key=_file_sort_key, reverse=True)
        if not candidates:
            raise HTTPException(404, f"目录 {path_str} 中未找到 .xlsx 文件")
        return candidates[0]

    raise HTTPException(400, f"无法识别路径：{path_str}")


def _list_report_files(path_str: str) -> List[tuple]:
    """列出报表目录下每天取用的 xlsx，返回 [(date_str, Path)] 升序。只列文件不解析内容。

    兼容日期子目录（YYYY-MM-DD/）与平铺（*_YYYYMMDD.xlsx）两种结构；供趋势增量入库用。
    """
    p = pathlib.Path(path_str)
    if p.is_file():
        p = p.parent
    if not p.is_dir():
        raise HTTPException(404, f"目录不存在：{path_str}")

    date_dirs = _list_date_dirs(p)
    file_list: List[tuple] = []
    if date_dirs:
        for d in reversed(date_dirs):  # 升序
            xlsxes = sorted(d.glob("*.xlsx"), key=lambda f: f.name, reverse=True)
            if xlsxes:
                file_list.append((d.name, xlsxes[0]))
    else:
        flat = sorted(
            (f for f in p.glob("*.xlsx") if _DATE_PAT.search(f.name)),
            key=_file_sort_key,
        )
        for f in flat:
            m = _DATE_PAT.search(f.name)
            ds = m.group(1)
            file_list.append((f"{ds[:4]}-{ds[4:6]}-{ds[6:]}", f))
    return file_list


# ─── PPT builder ────────────────────────────────────────────────────────────

def _build_pptx(data: dict) -> io.BytesIO:
    """缺陷统计报表 PPT：封面 + 三张矩阵表。

    表格排版全部走 pptx_utils 的统一入口（分页 / 列宽归一化 / 合计行加粗）——
    原来这里自己算行高与列宽：行一多表格顺着页面往下长到幻灯片外面，
    列一多每列被压到 0.4"，导出来没法直接用。
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    import pptx_utils as PU

    prs = PU._new_pres()
    blank = prs.slide_layouts[6]

    C_CRIT   = RGBColor(0x8E, 0x24, 0xAA)   # 致命
    C_RED    = RGBColor(0xC6, 0x28, 0x28)   # 严重
    C_ORANGE = RGBColor(0xE6, 0xA2, 0x3C)   # 一般
    C_GRAY   = RGBColor(0x90, 0x93, 0x99)   # 提示
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

    raw = data.get("raw", [])
    stamp = f"{data.get('actual_file', '')}   {data.get('file_mtime', '')}".strip()
    subtitle = f"数据源：{stamp}" if stamp else "缺陷统计报表"

    def _txt(slide, text, x, y, w, h, size=12, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        p = tb.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        PU._apply_run_font(run, size, bold, color)

    # ── 封面 ────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    _txt(cover, "缺陷统计报表", 1, 1.7, 11.3, 1.4, size=44, bold=True,
         color=PU._BRAND, align=PP_ALIGN.CENTER)
    _txt(cover, stamp, 1, 3.2, 11.3, 0.5, size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    def _sev(name):
        return sum(1 for r in raw if (r.get("severity") or "").strip() == name)

    cards = [("合计", len(raw), PU._BRAND)]
    for name, color in (("致命", C_CRIT), ("严重", C_RED), ("一般", C_ORANGE), ("提示", C_GRAY)):
        n = _sev(name)
        # 致命一栏为 0 时不铺出来：多数项目没有致命单，占着一张卡片只会稀释其它数字
        if n or name != "致命":
            cards.append((name, n, color))

    # 卡片整体居中：原来从 0.7" 起固定步长，卡片数一变就偏到一边
    card_w, gap = 2.5, 0.28
    total_w = len(cards) * card_w + (len(cards) - 1) * gap
    cx = (13.333 - total_w) / 2
    for label, val, clr in cards:
        shp = cover.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(cx), Inches(4.3), Inches(card_w), Inches(1.6))
        shp.fill.solid()
        shp.fill.fore_color.rgb = clr
        shp.line.fill.background()
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = False
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = str(val)
        PU._apply_run_font(r1, 36, True, C_WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        PU._apply_run_font(r2, 13, False, C_WHITE)
        cx += card_w + gap

    # 封面也带页脚：模板里页脚是每页固定件，只有封面没有的话，
    # 一叠幻灯片翻过去第一页会显得像另一份材料
    PU._add_footer(cover, 1, 1)

    # ── 三张矩阵表：列多了按列切页，行多了按行切页 ──────────────────
    for key, title in (("monthly_by_group", "按小组月度统计"),
                       ("by_customer", "按客户分布"),
                       ("feature_by_group", "特性 × 小组分布")):
        blk = data.get(key) or {}
        cols = blk.get("columns") or []
        rows = blk.get("rows") or []
        if not rows:
            continue
        PU.add_matrix_slides(
            prs, title, subtitle, "小组", cols,
            [[r["label"]] + [str(r.get(c, 0)) for c in cols] for r in rows],
        )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─── endpoints ──────────────────────────────────────────────────────────────

@router.get("/dates")
def list_dates(_: models.User = Depends(get_current_user)):
    """列出报表目录下所有含 xlsx 文件的日期子目录（YYYY-MM-DD），最新在前。"""
    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        return []
    p = pathlib.Path(path_str)
    if not p.is_dir():
        return []
    return [d.name for d in _list_date_dirs(p)]


@router.get("/data")
def get_data(date: Optional[str] = None, _: models.User = Depends(get_current_user)):
    """加载单张报表。date 参数为 YYYY-MM-DD，不传则取最新一天。"""
    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        return {"configured": False}
    target = _resolve_for_date(path_str, date)
    # 浅拷贝：缓存里存的是共享 dict，避免调用方补字段污染缓存
    result = dict(_parse_excel_cached(str(target)))
    result["actual_file"] = target.name
    result["date_dir"] = target.parent.name if _DATE_DIR_PAT.match(target.parent.name) else None
    return {"configured": True, **result}


@router.get("/trend")
def get_trend(db: Session = Depends(get_db), _: models.User = Depends(get_current_user)):
    """按天聚合的问题单趋势。数字入库（issue_report_daily），趋势直接读库；
    只有新增/变更（file_mtime 变化）的那天才重新解析对应 xlsx——增量、免全量重扫。
    """
    import openpyxl

    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        raise HTTPException(400, "未配置报表路径")

    file_list = _list_report_files(path_str)
    if not file_list:
        raise HTTPException(404, "目录中无可用报表（日期子目录或含日期后缀的 xlsx）")

    existing = {d.date: d for d in db.query(models.IssueReportDaily).all()}
    disk_dates = {ds for ds, _ in file_list}

    # 增量入库：只解析库里没有、或文件已变（名字/mtime 变化）的那些天
    for date_str, f in file_list:
        try:
            mtime = str(f.stat().st_mtime)
        except OSError:
            mtime = ""
        cur = existing.get(date_str)
        if cur is not None and cur.file_name == f.name and cur.file_mtime == mtime:
            continue  # 未变，跳过解析
        try:
            wb = openpyxl.load_workbook(str(f), data_only=True)
            raw = _parse_raw_from_wb(wb)
            wb.close()
        except Exception:
            continue
        if cur is None:
            cur = models.IssueReportDaily(date=date_str)
            db.add(cur)
            existing[date_str] = cur
        cur.file_name = f.name
        cur.file_mtime = mtime
        cur.total = len(raw)
        cur.by_group_json = json.dumps(_count_by(raw, "group"), ensure_ascii=False)
        cur.by_severity_json = json.dumps(_count_by(raw, "severity"), ensure_ascii=False)
        cur.ingested_at = datetime.utcnow()

    # 磁盘上已删除的日期：清掉库里的陈旧行，保持趋势与目录一致
    stale = [d for date, d in existing.items() if date not in disk_dates]
    for d in stale:
        db.delete(d)
    db.commit()

    rows = (
        db.query(models.IssueReportDaily)
        .order_by(models.IssueReportDaily.date.asc())
        .all()
    )
    daily: List[Dict] = []
    all_groups: set = set()
    all_severities: set = set()
    for r in rows:
        try:
            bg = json.loads(r.by_group_json or "{}")
            bs = json.loads(r.by_severity_json or "{}")
        except (ValueError, TypeError):
            bg, bs = {}, {}
        all_groups.update(bg.keys())
        all_severities.update(bs.keys())
        daily.append({
            "date": r.date, "file": r.file_name,
            "total": r.total, "by_group": bg, "by_severity": bs,
        })

    sev_order = ["严重", "一般", "提示"]
    return {
        "daily":          daily,
        "all_groups":     sorted(all_groups),
        "all_severities": sorted(all_severities, key=lambda s: sev_order.index(s) if s in sev_order else 99),
    }


# ─── 通过脚本调用外部 API 拉取问题单（按项目）──────────────────────────────────
def _normalize_issue_row(r: dict) -> Dict[str, str]:
    """把脚本返回的一条问题单规整成「原始数据」表同款字段（缺的留空）。"""
    return {col: (str(r.get(col)).strip() if r.get(col) is not None else "") for col in _RAW_COLS}


def _script_timeout(cfg: Dict) -> int:
    """脚本执行超时（秒）。config.issue_script_timeout，默认 600s。

    原来写死 120s：DTS 接口慢一点就被杀，但快照有时已经落盘，表现为"报错了、
    过会儿刷新数据又有了"。留成可配置，默认给足。
    """
    try:
        v = int(cfg.get("issue_script_timeout") or 0)
    except (TypeError, ValueError):
        v = 0
    return v if v > 0 else 600


def _run_issue_api_script(project: str) -> List[Dict]:
    """以 `python <issue_api_script_path> <project>` 调用脚本，期望 stdout 为 JSON 数组。"""
    cfg = _load_config()
    script = (cfg.get("issue_api_script_path") or "").strip()
    if not script:
        raise HTTPException(400, "未配置 API 脚本（issue_api_script_path）")
    sp = pathlib.Path(script)
    if not sp.exists():
        raise HTTPException(404, f"脚本不存在：{script}")

    timeout = _script_timeout(cfg)
    cmd = [sys.executable, str(sp), project] if sp.suffix.lower() == ".py" else [str(sp), project]
    try:
        result = subprocess.run(
            cmd, capture_output=True, encoding="utf-8", errors="replace",
            timeout=timeout, cwd=str(sp.parent),
        )
    except subprocess.TimeoutExpired:
        raise HTTPException(504, f"脚本执行超时（>{timeout} 秒），可在「配置」中调大超时时间")
    except Exception as exc:
        raise HTTPException(500, f"脚本启动失败：{exc}")

    if result.returncode != 0:
        raise HTTPException(500, f"脚本退出码 {result.returncode}：{(result.stderr or '')[-500:]}")
    out = (result.stdout or "").strip()
    if not out:
        raise HTTPException(500, "脚本无输出（应向 stdout 打印 JSON 数组）")
    try:
        data = json.loads(out)
    except json.JSONDecodeError as exc:
        raise HTTPException(500, f"脚本输出不是合法 JSON：{exc}")
    if not isinstance(data, list):
        raise HTTPException(500, "脚本输出应为问题单数组（JSON list）")
    return [_normalize_issue_row(r) for r in data if isinstance(r, dict)]


# ─── 每日快照：库存"数字"（趋势）+ 文件存明细（钻取）──────────────────────────
_BACKEND_DIR = pathlib.Path(__file__).resolve().parent.parent
_SEV_ORDER = {"严重": 0, "一般": 1, "提示": 2}


def _snapshot_root() -> pathlib.Path:
    """快照明细文件根目录：优先 config.issue_snapshot_dir，否则 backend/data/issue_snapshots。"""
    cfg = _load_config()
    d = (cfg.get("issue_snapshot_dir") or "").strip()
    root = pathlib.Path(d) if d else (_BACKEND_DIR / "data" / "issue_snapshots")
    root.mkdir(parents=True, exist_ok=True)
    return root


def _safe_slug(s: str) -> str:
    return re.sub(r"[^\w\-]", "_", s or "") or "_"


# ─── 采集后富化：部门过滤 + 责任人归组 + 从标题提取客户面 ─────────────────────
def _as_str_list(v) -> List[str]:
    """config 里可能存成 list 或分号/换行分隔的字符串，统一成去空的列表。"""
    if isinstance(v, list):
        return [str(x).strip() for x in v if str(x).strip()]
    if isinstance(v, str):
        return [s.strip() for s in re.split(r"[;；\n]", v) if s.strip()]
    return []


def _load_issue_groups(cfg: Dict) -> List[tuple]:
    """config.issue_groups: [{name, members}] → [(小组名, [成员,...])]。成员支持分号/换行分隔。"""
    groups: List[tuple] = []
    for g in (cfg.get("issue_groups") or []):
        if not isinstance(g, dict):
            continue
        name = str(g.get("name") or "").strip()
        if not name:
            continue
        groups.append((name, _as_str_list(g.get("members"))))
    return groups


def _name_class(ch: str) -> str:
    """字符在「名字」意义上的类别：分隔符 / 汉字 / 西文（字母数字算一类）。

    用来判断一次子串匹配是不是把一个更长的名字从中间切开了。
    """
    if not ch.isalnum():
        return ""            # 空格、括号、斜杠、逗号……都算分隔符
    return "cjk" if "\u4e00" <= ch <= "\u9fff" else "latin"


def _contains_person_name(text: str, name: str) -> bool:
    """`name` 是否作为一个**完整的人名**出现在 `text` 里。

    归组必须做模糊匹配：DTS 的责任人字段常常是「张伟 00123456」「张伟(zhangwei)」
    这种姓名带工号/英文名的写法，而配置里的小组名单只写姓名。但**朴素的子串包含**
    会让「张伟」认走「张伟明」的单——名单里没有张伟明，他的单却被安安静静地记到了
    张伟所在的组，组级负载和交叉表都因此偏一点，而两边看着都对。

    所以要求匹配处的两端是**真正的边界**：紧邻的字符要么是分隔符，要么属于另一个
    字符类。汉字接汉字（张伟|明）＝切开了一个更长的名字，拒绝；汉字接西文
    （张伟|00123456、张伟|zhangwei）是姓名与工号/英文名的交界，接受。

    注意这条规则**比客户面的 `_contains_name()` 严**，两者不能互相套用：客户名是从
    问题单**标题**里认的，标题是自然语句、名字两侧本来就没有分隔符（「西安1号机异常」），
    按这里的规则会被整批拒掉；而责任人是一个**名字字段**，两端本就该是边界。
    """
    if not name:
        return False
    head, tail = _name_class(name[0]), _name_class(name[-1])
    start = 0
    while True:
        i = text.find(name, start)
        if i < 0:
            return False
        end = i + len(name)
        head_ok = not head or i == 0 or _name_class(text[i - 1]) != head
        tail_ok = not tail or end >= len(text) or _name_class(text[end]) != tail
        if head_ok and tail_ok:
            return True
        start = i + 1


def _match_group(owner: str, groups: List[tuple]) -> str:
    """按责任人姓名匹配所属小组（大小写不敏感 + 带边界的双向包含匹配）。

    双向是有意的：名单可能写得比 DTS 短（只写姓名，DTS 里带工号），也可能写得更长
    （名单里写「张伟(SE)」，DTS 里只有姓名）。两个方向都过 `_contains_person_name()`
    的边界检查，所以「张伟」不会再认走「张伟明」。
    认不上的人不丢行，由调用方归到 UNGROUPED_GROUP 并报给管理员去补名单。
    """
    o = (owner or "").strip().lower()
    if not o:
        return ""
    for name, members in groups:
        for m in members:
            ml = (m or "").strip().lower()
            if not ml:
                continue
            if ml == o or _contains_person_name(o, ml) or _contains_person_name(ml, o):
                return name
    return ""


def _load_customer_matchers(db: Session) -> List[tuple]:
    """从客户主数据（code/全称/别名）构建 [(匹配文本_lower, 展示名)]，按长度降序（优先更具体）。

    长度降序只解决「两个名字都登记了」的情形（「11号机」比「1号机」长，先试到先命中）。
    名字被另一个名字从中间切开的那一类要靠 `_contains_name()` 的数字边界挡——两道缺一不可。
    """
    matchers: List[tuple] = []
    customers = db.query(models.Customer).filter(models.Customer.is_active == True).all()  # noqa: E712
    id2label = {}
    for c in customers:
        label = (c.display_name or c.code or "").strip()
        id2label[c.id] = label
        for t in (c.code, c.display_name):
            if t and t.strip():
                matchers.append((t.strip().lower(), label))
    for a in db.query(models.CustomerAlias).all():
        label = id2label.get(a.customer_id)
        if label and a.alias and a.alias.strip():
            matchers.append((a.alias.strip().lower(), label))
    seen, uniq = set(), []
    for mt, label in matchers:
        if mt and mt not in seen:
            seen.add(mt)
            uniq.append((mt, label))
    uniq.sort(key=lambda x: len(x[0]), reverse=True)
    return uniq


def _contains_name(text: str, name: str) -> bool:
    """`name` 是否作为一个完整的名字出现在 `text` 里——**不接受把一串数字从中间切开**。

    客户面在这套数据里大量是「N号机」这种带编号的名字，而 `"1号机" in "11号机"` 是真的
    （从第二个字符起就是）。于是「11号机」的单子会落到「1号机」那一档：两台机器的问题
    混进同一行，同一台机器的问题散在两行——数字还是那些数字、加起来也对得上，
    没人会把它当成 bug 报上来。

    按长度降序试（`_load_customer_matchers`）只挡得住「两台机器都在客户主数据里」的情形；
    只要「11号机」没登记、或写法对不上，「1号机」照样会把它吃掉。所以这里再加一道边界：
    匹配文本以数字开头时，它前面一位不能还是数字；以数字结尾时，后面一位不能还是数字。
    只拒绝「数字被切开」这一种，中文/英文边界不管——客户名混在标题里本来就没有分隔符，
    要求两侧都是分隔符会把「西安1号机异常」这类正常标题一并拒掉。
    """
    if not name:
        return False
    start = 0
    while True:
        i = text.find(name, start)
        if i < 0:
            return False
        end = i + len(name)
        head_ok = not (name[0].isdigit() and i > 0 and text[i - 1].isdigit())
        tail_ok = not (name[-1].isdigit() and end < len(text) and text[end].isdigit())
        if head_ok and tail_ok:
            return True
        start = i + 1


def _match_customer(title: str, matchers: List[tuple]) -> str:
    t = (title or "").lower()
    if not t:
        return ""
    for mt, label in matchers:
        if _contains_name(t, mt):
            return label
    return ""


def _enrich_rows(db: Session, rows: List[Dict]) -> List[Dict]:
    """对采集到的问题单做：① 部门过滤 ② 按责任人归组 ③ 从标题提取客户面。

    配置项（config.json，问题单管理「配置」tab 维护）：
      issue_exclude_statuses —— 直接剔除的状态（默认 关闭/撤销；子串匹配进展/状态）
      issue_stat_departments —— 只统计这些部门（子串匹配责任人部门全路径；留空＝全部）
      issue_groups           —— [{name, members}]，成员分号分隔，按责任人归组
    客户面来自客户主数据（客户面管理），用 code/全称/别名 在标题里做包含匹配。

    **部门过滤会丢行，小组归组不会**：部门答的是"这单归不归我们管"，答否就该出统计；
    小组答的是"归我们哪个组"，答不上来只说明名单没维护全，不是这单不该统计。
    """
    cfg = _load_config()
    exclude = _as_str_list(cfg.get("issue_exclude_statuses")) or ["关闭", "撤销"]
    depts = _as_str_list(cfg.get("issue_stat_departments"))
    groups = _load_issue_groups(cfg)
    matchers = _load_customer_matchers(db)

    out: List[Dict] = []
    for r in rows:
        # ① 先剔除已关闭 / 已撤销 的单（任何地方都不出现）
        prog = r.get("progress", "") or ""
        if any(s and s in prog for s in exclude):
            continue
        # ② 部门过滤：匹配责任人部门全路径（兼容部门落在上级字段的情况），回退直属部门
        if depts:
            dept = (r.get("dept_path") or "") or (r.get("department") or "")
            if not any(d in dept for d in depts):
                continue
        # ③ 按责任人归组。**不在名单里的人不丢**，归到「未归组」。
        # 问题单从定位转到实施修改时换责任人是正常流转，新人/借调/换部门的人
        # 不在名单里很常见；丢掉这一行的后果不是"少统计一条"，而是它在下一次
        # 相邻快照差分里凭空变成一笔「解决」（见 _ensure_flows：解决＝上次有、
        # 今天没有，从不读状态）。所以留下来，并由 _ungrouped_owners() 报出去，
        # 让管理员去补名单——名单不全是配置问题，不该表现成数据问题。
        if groups:
            r["group"] = _match_group(r.get("owner", ""), groups) or UNGROUPED_GROUP
        # ④ 从标题提取客户面
        if matchers and not (r.get("customer") or "").strip():
            r["customer"] = _match_customer(r.get("title", ""), matchers)
        out.append(r)
    return out


def _ungrouped_owners(rows: List[Dict]) -> List[Dict]:
    """本次采集里归不到小组的责任人：[{owner, dept, count}]，按条数降序。

    这不是统计口径，是**待办清单**——每一条都对应「配置里的小组名单少了一个人」。
    带上部门是为了让管理员一眼看出该往哪个组里加，而不用回 DTS 里查这人是谁。
    """
    acc: Dict[str, Dict] = {}
    for r in rows:
        if (r.get("group") or "") != UNGROUPED_GROUP:
            continue
        owner = str(r.get("owner") or "").strip() or "（无责任人）"
        item = acc.get(owner)
        if item is None:
            item = acc[owner] = {"owner": owner, "dept": "", "count": 0}
        item["count"] += 1
        if not item["dept"]:
            item["dept"] = (r.get("department") or r.get("dept_path") or "").strip()
    return sorted(acc.values(), key=lambda x: (-x["count"], x["owner"]))


def _take_snapshot(db: Session, project: str,
                   source: str = "api") -> tuple[models.IssueSnapshot, List[Dict]]:
    """拉取该项目问题单 → 明细写文件、聚合数字写库（同项目同日覆盖）。

    返回 (快照, 归不到小组的责任人清单)。后者是给管理员看的待办，不入库——
    它完全由明细文件推得（`/ungrouped` 随时重算），存一份就得考虑何时失效。

    可能抛 HTTPException（脚本未配置 / 执行失败）——调用方按需捕获。
    """
    raw = _run_issue_api_script(project)
    raw = _enrich_rows(db, raw)   # 部门过滤 + 责任人归组 + 标题提取客户面
    today = datetime.now().strftime("%Y-%m-%d")

    # 明细落文件（<项目>/<日期>.json）
    root = _snapshot_root()
    rel = f"{_safe_slug(project)}/{today}.json"
    fp = root / rel
    fp.parent.mkdir(parents=True, exist_ok=True)
    fp.write_text(json.dumps(raw, ensure_ascii=False), encoding="utf-8")

    # 自动落盘 Excel 备份：原始表 + 分析表（同日多次用时间戳不覆盖）
    _export_snapshot_excel(project, raw, today)

    # upsert 快照元数据
    snap = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project,
                models.IssueSnapshot.snapshot_date == today)
        .first()
    )
    if snap is None:
        snap = models.IssueSnapshot(project=project, snapshot_date=today)
        db.add(snap)
    snap.total = len(raw)
    snap.data_file = rel
    snap.source = source
    snap.created_at = datetime.utcnow()
    db.flush()  # 拿到 snap.id

    # 重建维度聚合数字（group / customer / severity）
    db.query(models.IssueSnapshotStat).filter(
        models.IssueSnapshotStat.snapshot_id == snap.id
    ).delete(synchronize_session=False)
    for dim in ("group", "customer", "severity"):
        for key, cnt in _count_by(raw, dim).items():
            db.add(models.IssueSnapshotStat(
                snapshot_id=snap.id, dimension=dim, dim_key=key, count=cnt,
            ))
    db.commit()
    db.refresh(snap)

    # 顺手算当天的新增/解决差分（只比对上一次快照，看图时就不用现读一堆文件了）
    try:
        _ensure_flows(db, project)
    except Exception:
        db.rollback()   # 差分是附加信息，算不出来不影响这次采集
    return snap, _ungrouped_owners(raw)


def collect_with_log(db: Session, project: str, source: str = "auto") -> Dict[str, Any]:
    """采集一个项目并写执行日志（成功失败都写）。不抛异常，返回结果字典。

    定时任务与手动采集共用，保证两条路径的日志口径一致。
    """
    started = datetime.utcnow()
    log = models.IssueCollectLog(project=project, source=source, started_at=started)
    result: Dict[str, Any]
    try:
        snap, ungrouped = _take_snapshot(db, project, source=source)
        log.ok, log.total, log.error = True, snap.total, ""
        # ungrouped 随采集结果一并回给页面：名单少人是配置问题，采集当场提示
        # 比等人自己去翻配置页有用得多。
        result = {"project": project, "ok": True, "date": snap.snapshot_date,
                  "total": snap.total, "ungrouped": ungrouped}
    except HTTPException as exc:
        log.ok, log.error = False, str(exc.detail)[:2000]
        result = {"project": project, "ok": False, "error": log.error}
    except Exception as exc:  # noqa: BLE001 —— 脚本/解析的意外错误也要留痕
        log.ok, log.error = False, f"{type(exc).__name__}: {exc}"[:2000]
        result = {"project": project, "ok": False, "error": log.error}

    log.finished_at = datetime.utcnow()
    log.duration_ms = int((log.finished_at - started).total_seconds() * 1000)
    try:
        db.add(log)
        db.commit()
    except Exception:  # 日志写失败不能影响采集结论
        db.rollback()
    return result


# ─── 手动采集：后台线程执行 + 轮询状态 ────────────────────────────────────────
# 采集要跑几分钟，而前端 axios 超时只有 10s：同步返回必然先弹「采集失败」、
# 但后台其实还在跑并最终成功，于是"过会儿刷新数据又有了"。改成立即返回 + 轮询。
_collect_lock = threading.Lock()
_collect_state: Dict[str, Any] = {
    "running": False, "projects": [], "current": None,
    "started_at": None, "finished_at": None, "results": [],
}


def _collect_worker(projects: List[str], source: str) -> None:
    db = SessionLocal()
    try:
        for p in projects:
            _collect_state["current"] = p
            _collect_state["results"].append(collect_with_log(db, p, source=source))
    finally:
        db.close()
        _collect_state["current"] = None
        _collect_state["finished_at"] = iso_local(datetime.utcnow())
        _collect_state["running"] = False
        if _collect_lock.locked():
            _collect_lock.release()


@router.post("/snapshot-collect")
def snapshot_collect(
    request: Request,
    project: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(require_admin),
):
    """启动一次采集（仅管理员），立即返回；用 GET /collect-status 轮询进度。

    project 省略则采集 config.issue_api_projects 全部。
    """
    cfg = _load_config()
    projects = [project] if project else (cfg.get("issue_api_projects") or [])
    if not projects:
        raise HTTPException(400, "没有可采集的项目（未配置 issue_api_projects）")

    if not _collect_lock.acquire(blocking=False):
        raise HTTPException(423, f"已有采集任务在执行（{_collect_state.get('current') or '…'}），请稍候")

    _collect_state.update({
        "running": True, "projects": list(projects), "current": None,
        "started_at": iso_local(datetime.utcnow()), "finished_at": None, "results": [],
    })
    threading.Thread(target=_collect_worker, args=(list(projects), "manual"), daemon=True).start()

    log_op(db, action="issue_snapshot", target="issue_snapshot", target_id=None,
           detail=f"手动触发采集：{', '.join(projects)}", user=current_user, request=request)
    return {"started": True, "projects": projects}


@router.get("/collect-status")
def collect_status(_: models.User = Depends(get_current_user)):
    """采集任务进度（所有登录用户可查）。running=false 且 results 非空即为本轮结果。"""
    return dict(_collect_state)


@router.get("/collect-schedule")
def collect_schedule(_: models.User = Depends(get_current_user)):
    """定时采集的运行态（所有登录用户可查）：是否启用、下次触发时间、是否配全。

    读权限不收紧到 admin：普通用户看到"下次采集 明天 07:30"就不会来问数据什么时候更新。
    字段含义与排障口径见 scheduler.snapshot_job_status()。
    """
    try:
        import scheduler
        return scheduler.snapshot_job_status()
    except Exception as exc:  # 调度模块装载失败也要给前端一个明确结论
        return {"scheduler_running": False, "enabled": False, "time": "",
                "job_registered": False, "next_run_at": None,
                "script_path": "", "projects": [], "error": f"{type(exc).__name__}: {exc}"}


@router.get("/collect-logs")
def collect_logs(project: Optional[str] = None, limit: int = 50,
                 db: Session = Depends(get_db),
                 _: models.User = Depends(get_current_user)):
    """采集执行日志（新→旧）。project 省略则返回所有项目。"""
    q = db.query(models.IssueCollectLog)
    if project:
        q = q.filter(models.IssueCollectLog.project == project)
    rows = q.order_by(models.IssueCollectLog.started_at.desc()).limit(max(1, min(limit, 200))).all()
    return [
        {
            "id": r.id, "project": r.project, "source": r.source, "ok": bool(r.ok),
            "total": r.total or 0, "duration_ms": r.duration_ms or 0, "error": r.error or "",
            # 库里存的是 UTC（datetime.utcnow），出口转北京时间，否则页面上
            # 07:30 跑的定时采集会显示成头一天 23:30
            "started_at": fmt_local(r.started_at),
            "finished_at": fmt_local(r.finished_at),
        }
        for r in rows
    ]


@router.get("/snapshots")
def list_snapshots(project: str, db: Session = Depends(get_db),
                   _: models.User = Depends(get_current_user)):
    """某项目的快照列表（新→旧），只含元数据数字。"""
    rows = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project)
        .order_by(models.IssueSnapshot.snapshot_date.desc())
        .all()
    )
    return [
        {"id": r.id, "date": r.snapshot_date, "total": r.total, "source": r.source,
         "created_at": fmt_local(r.created_at)}
        for r in rows
    ]


@router.get("/snapshot-detail")
def snapshot_detail(project: str, date: Optional[str] = None,
                    db: Session = Depends(get_db),
                    _: models.User = Depends(get_current_user)):
    """某次快照的明细（从文件加载完整行）；date 省略取最新。"""
    q = db.query(models.IssueSnapshot).filter(models.IssueSnapshot.project == project)
    snap = (q.filter(models.IssueSnapshot.snapshot_date == date).first() if date
            else q.order_by(models.IssueSnapshot.snapshot_date.desc()).first())
    if snap is None:
        return {"exists": False, "project": project}
    raw: List[Dict] = []
    try:
        fp = _snapshot_root() / snap.data_file
        if fp.exists():
            raw = json.loads(fp.read_text(encoding="utf-8"))
    except Exception:
        raw = []
    return {
        "exists": True, "project": project, "date": snap.snapshot_date,
        "created_at": fmt_local(snap.created_at),
        "total": snap.total, "count": len(raw), "raw": raw,
        "by_severity": _count_by(raw, "severity"),
        "by_group": _count_by(raw, "group"),
        "by_customer": _count_by(raw, "customer"),
    }


@router.get("/snapshot-trend")
def snapshot_trend(project: str, dimension: str = "group",
                   db: Session = Depends(get_db),
                   _: models.User = Depends(get_current_user)):
    """趋势：只从库里读维度聚合数字（不碰明细文件）。dimension ∈ group/customer/severity。"""
    if dimension not in ("group", "customer", "severity"):
        dimension = "group"
    snaps = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project)
        .order_by(models.IssueSnapshot.snapshot_date.asc())
        .all()
    )
    if not snaps:
        return {"project": project, "dimension": dimension, "dates": [], "total": [], "series": []}

    dates = [s.snapshot_date for s in snaps]
    total = [s.total for s in snaps]
    id_to_idx = {s.id: i for i, s in enumerate(snaps)}
    stats = (
        db.query(models.IssueSnapshotStat)
        .filter(models.IssueSnapshotStat.dimension == dimension,
                models.IssueSnapshotStat.snapshot_id.in_([s.id for s in snaps]))
        .all()
    )
    matrix: Dict[str, List[int]] = {}
    order: List[str] = []
    for st in stats:
        idx = id_to_idx.get(st.snapshot_id)
        if idx is None:
            continue
        if st.dim_key not in matrix:
            matrix[st.dim_key] = [0] * len(dates)
            order.append(st.dim_key)
        matrix[st.dim_key][idx] = st.count
    if dimension == "severity":
        order.sort(key=lambda k: _SEV_ORDER.get(k, 99))
    else:
        order.sort()
    series = [{"name": k, "data": matrix[k]} for k in order]
    return {"project": project, "dimension": dimension, "dates": dates,
            "total": total, "series": series}


# ─── 快照导出 Excel：原始数据 + 统计分析 两张表 ────────────────────────────────
def _cross_table(rows: List[Dict], row_field: str, col_field: str,
                 col_order: Optional[List[str]] = None,
                 row_fallback: str = "未标注") -> Dict[str, Any]:
    """行维度 × 列维度 交叉计数 → {columns:[...,'合计'], rows:[{label,...}], total_row}。"""
    matrix: Dict[str, Dict[str, int]] = {}
    col_totals: Dict[str, int] = {}
    col_seen: List[str] = []
    grand = 0
    for r in rows:
        rv = (r.get(row_field) or "").strip() or row_fallback
        cv = (r.get(col_field) or "").strip() or "未标注"
        if cv not in col_totals:
            col_totals[cv] = 0
            col_seen.append(cv)
        matrix.setdefault(rv, {})
        matrix[rv][cv] = matrix[rv].get(cv, 0) + 1
        col_totals[cv] += 1
        grand += 1
    if col_order:
        cols = [c for c in col_order if c in col_totals] + sorted(c for c in col_seen if c not in col_order)
    else:
        cols = sorted(col_seen)
    columns = cols + ["合计"]
    out_rows = []
    for rv in sorted(matrix.keys()):
        rec: Dict[str, Any] = {"label": rv}
        t = 0
        for c in cols:
            rec[c] = matrix[rv].get(c, 0)
            t += rec[c]
        rec["合计"] = t
        out_rows.append(rec)
    total_row = {"label": "合计", **{c: col_totals.get(c, 0) for c in cols}, "合计": grand}
    return {"columns": columns, "rows": out_rows, "total_row": total_row}


_RAW_XLSX_COLS = [
    ("issue_id", "缺陷业务编号", 20), ("title", "标题", 42), ("owner", "当前责任人", 12),
    ("group", "所属小组", 14), ("department", "责任人部门", 20), ("customer", "客户面", 14),
    ("feature", "特性", 14), ("subsystem", "子系统", 14), ("module", "模块", 14),
    ("progress", "进展", 12), ("severity", "严重程度", 10), ("year_month", "年月", 10),
    ("version", "版本信息", 22),
]


def _fill_raw_sheet(ws, raw: List[Dict]) -> None:
    """原始数据表：问题单明细，一行一条。"""
    from xlsx_io import beautify, style_header

    ws.title = "原始数据"
    style_header(ws, [h for _, h, _ in _RAW_XLSX_COLS])
    for c_idx, (_, _, w) in enumerate(_RAW_XLSX_COLS, start=1):
        ws.column_dimensions[ws.cell(1, c_idx).column_letter].width = w
    for r in raw:
        ws.append([r.get(k, "") for k, _, _ in _RAW_XLSX_COLS])
    # 短列（责任人/小组/进展/严重程度/年月等）居中
    beautify(ws, center_cols={3, 4, 6, 10, 11, 12})


def _fill_analysis_sheet(ws, raw: List[Dict]) -> None:
    """统计分析表：按小组 / 客户面 / 研发问题 / 年月 × 严重程度 四张交叉表纵向排布。

    客户面表只统计标题匹配到客户的单子；匹配不到的是研发问题，单独一张按小组的表
    （口径与前端 IssueApiPanel 一致，页面与导出必须同款）。
    """
    import brand
    from openpyxl.styles import Alignment, Font, PatternFill
    # 配色走 brand.py，与清单类导出和 PPT 同一套；这里曾经自己写死品牌蓝
    head_font = Font(bold=True, color=brand.HEADER_TEXT)
    head_fill = PatternFill("solid", fgColor=brand.HEADER_BG)
    title_font = Font(bold=True, size=12, color=brand.BRAND)
    total_font = Font(bold=True)
    center = Alignment(horizontal="center", vertical="center")
    ws.title = "统计分析"
    ws.column_dimensions["A"].width = 18
    SEV = ["严重", "一般", "提示"]
    row_ptr = [1]   # 显式维护当前写入行，避免依赖 max_row

    def _write_cross(row_label: str, title: str, cross: Dict[str, Any]):
        r = row_ptr[0]
        ws.cell(r, 1, title).font = title_font
        r += 1
        for c_idx, val in enumerate([row_label] + cross["columns"], start=1):
            cell = ws.cell(r, c_idx, val)
            cell.font = head_font
            cell.fill = head_fill
            cell.alignment = center
        for row in cross["rows"]:
            r += 1
            ws.cell(r, 1, row["label"])
            for c_idx, col in enumerate(cross["columns"], start=2):
                ws.cell(r, c_idx, row.get(col, 0)).alignment = center
        r += 1
        ws.cell(r, 1, cross["total_row"]["label"]).font = total_font
        for c_idx, col in enumerate(cross["columns"], start=2):
            cell = ws.cell(r, c_idx, cross["total_row"].get(col, 0))
            cell.font = total_font
            cell.alignment = center
        row_ptr[0] = r + 2   # 空一行再写下一张表

    cus_rows = [r for r in raw if (r.get("customer") or "").strip()]
    dev_rows = [r for r in raw if not (r.get("customer") or "").strip()]

    _write_cross("小组", "按小组 × 严重程度", _cross_table(raw, "group", "severity", SEV, UNGROUPED_GROUP))
    _write_cross("客户面", f"按客户面 × 严重程度（客户面问题 {len(cus_rows)} 条）",
                 _cross_table(cus_rows, "customer", "severity", SEV, "未标注"))
    _write_cross("小组", f"研发问题 × 严重程度（{len(dev_rows)} 条，标题未匹配到客户）",
                 _cross_table(dev_rows, "group", "severity", SEV, UNGROUPED_GROUP))
    _write_cross("年月", "按年月 × 严重程度", _cross_table(raw, "year_month", "severity", SEV, "未标注"))


def _excel_dir(which: str) -> pathlib.Path:
    """自动导出目录：which=raw（原始表）/ analysis（分析表），各自可配置，
    默认 backend/data/issue_excel/<which>。"""
    cfg = _load_config()
    key = "issue_excel_raw_dir" if which == "raw" else "issue_excel_analysis_dir"
    d = (cfg.get(key) or "").strip()
    root = pathlib.Path(d) if d else (_BACKEND_DIR / "data" / "issue_excel" / which)
    root.mkdir(parents=True, exist_ok=True)
    return root


def _export_snapshot_excel(project: str, raw: List[Dict], date_str: str) -> None:
    """采集后自动落盘：原始表 + 分析表 各存一份到备份目录（同日多次用时间戳，不覆盖）。

    失败只吞掉不影响采集主流程。
    """
    try:
        import openpyxl
        slug = _safe_slug(project)
        stem = f"{date_str}_{datetime.now().strftime('%H%M%S')}"

        raw_wb = openpyxl.Workbook()
        _fill_raw_sheet(raw_wb.active, raw)
        raw_dir = _excel_dir("raw") / slug
        raw_dir.mkdir(parents=True, exist_ok=True)
        raw_wb.save(str(raw_dir / f"{slug}_原始_{stem}.xlsx"))

        ana_wb = openpyxl.Workbook()
        _fill_analysis_sheet(ana_wb.active, raw)
        ana_dir = _excel_dir("analysis") / slug
        ana_dir.mkdir(parents=True, exist_ok=True)
        ana_wb.save(str(ana_dir / f"{slug}_分析_{stem}.xlsx"))
    except Exception:
        pass


# ─── 每日新增 / 解决：相邻快照差分 ──────────────────────────────────────────
# 快照存的是"当天还开着的单"（关闭/撤销在 _enrich_rows 里已被剔除），所以：
#   新增 = 今天有、上次没有；解决 = 上次有、今天没有。
# 注意这里**从不读状态**：一单只要从快照里消失就算「解决」。所以任何"因为过滤规则
# 而掉出快照"的行都会变成一笔假解决——责任人归组因此刻意不丢行（见 _enrich_rows ③）。
# 差分要读明细文件，因此结果落 issue_snapshot_flows，采集后增量算一次，看图只读数字。
_ISSUE_NO_DATE = re.compile(r"^[A-Za-z]*(\d{4})(\d{2})(\d{2})\d*$")


def _issue_no_date(issue_id: str) -> str:
    """从缺陷业务编号提取创建日：SDTS + YYYY + MM + DD + 序号 → 'YYYY-MM-DD'，取不到返回 ''。"""
    m = _ISSUE_NO_DATE.match((issue_id or "").strip())
    if not m:
        return ""
    y, mo, d = m.groups()
    if not ("2000" <= y <= "2099" and "01" <= mo <= "12" and "01" <= d <= "31"):
        return ""
    return f"{y}-{mo}-{d}"


def _snapshot_ids(snap: models.IssueSnapshot) -> Optional[set]:
    """读某次快照明细里的编号集合；文件缺失/损坏返回 None（该天整个跳过，不参与差分）。

    返回空集合与返回 None 是两回事：前者表示"那天真的一条都没有"，
    后者表示"这天的数据丢了"——把丢数据的一天当成 0 条，会凭空冒出一大批"解决"。
    """
    if not snap.data_file:
        return None
    try:
        fp = _snapshot_root() / snap.data_file
        if not fp.exists():
            return None
        rows = json.loads(fp.read_text(encoding="utf-8"))
    except Exception:
        return None
    if not isinstance(rows, list):
        return None
    return {str(r.get("issue_id") or "").strip() for r in rows if isinstance(r, dict)
            and str(r.get("issue_id") or "").strip()}


def _ensure_flows(db: Session, project: str) -> List[models.IssueSnapshotFlow]:
    """按日期顺序补齐该项目的差分行（已算过且比对的上一天没变的直接复用）。

    存量快照是在这张表之前采集的，第一次看图时在这里补算；之后每天采集只算新的一天。
    """
    snaps = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project)
        .order_by(models.IssueSnapshot.snapshot_date.asc())
        .all()
    )
    if not snaps:
        return []
    existing = {
        f.snapshot_id: f
        for f in db.query(models.IssueSnapshotFlow)
        .filter(models.IssueSnapshotFlow.snapshot_id.in_([s.id for s in snaps])).all()
    }
    cache: Dict[int, Optional[set]] = {}

    def ids_of(snap) -> Optional[set]:
        if snap.id not in cache:
            cache[snap.id] = _snapshot_ids(snap)
        return cache[snap.id]

    out: List[models.IssueSnapshotFlow] = []
    prev: Optional[models.IssueSnapshot] = None
    dirty = False
    for snap in snaps:
        expect_prev = prev.snapshot_date if prev else ""
        flow = existing.get(snap.id)
        if flow is None or (flow.prev_date or "") != expect_prev:
            cur = ids_of(snap)
            if cur is None:          # 这天的明细文件没了 —— 整天跳过，prev 保持不变
                continue
            if prev is None:
                created, resolved, baseline = cur, set(), True
            else:
                pids = ids_of(prev) or set()
                created, resolved, baseline = cur - pids, pids - cur, False
            if flow is None:
                flow = models.IssueSnapshotFlow(snapshot_id=snap.id)
                db.add(flow)
            flow.project = project
            flow.snapshot_date = snap.snapshot_date
            flow.prev_date = expect_prev
            flow.is_baseline = baseline
            flow.created_count = len(created)
            flow.resolved_count = len(resolved)
            flow.created_ids_json = json.dumps(sorted(created), ensure_ascii=False)
            flow.resolved_ids_json = json.dumps(sorted(resolved), ensure_ascii=False)
            flow.computed_at = datetime.utcnow()
            dirty = True
        # 已算过的天直接复用：不再去读它的明细文件，否则每次看图都要重读整个目录
        out.append(flow)
        prev = snap
    if dirty:
        db.commit()
    return out


@router.get("/snapshot-flow")
def snapshot_flow(project: str, db: Session = Depends(get_db),
                  _: models.User = Depends(get_current_user)):
    """每日新增 / 解决 / 存量。两套口径一起返回，前端切换：

    - by_snapshot：按采集日差分。新增与解决同口径，净增＝新增−解决＝存量差，图能自洽；
      只能从第二次快照算起（首次是基线）。
    - by_issue_no：按缺陷编号里的创建日（SDTS+YYYYMMDD）统计新增。能回溯到开始采集之前，
      但只覆盖"被任何一次快照见过"的单——首次采集前就已闭环的单，谁也看不见了。
    """
    flows = _ensure_flows(db, project)
    snaps = {
        s.snapshot_date: s
        for s in db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project).all()
    }
    dates, created, resolved, open_cnt, net = [], [], [], [], []
    no_hist: Dict[str, int] = {}
    unknown_no = 0
    baseline_date = ""
    for f in flows:
        if f.is_baseline:
            baseline_date = f.snapshot_date
        else:
            dates.append(f.snapshot_date)
            created.append(f.created_count)
            resolved.append(f.resolved_count)
            snap = snaps.get(f.snapshot_date)
            open_cnt.append(snap.total if snap else 0)
            net.append(f.created_count - f.resolved_count)
        # 基线那天的编号也要计入"按编号日期"曲线：它们同样是历史上某天新增的
        for iid in json.loads(f.created_ids_json or "[]"):
            d = _issue_no_date(iid)
            if d:
                no_hist[d] = no_hist.get(d, 0) + 1
            else:
                unknown_no += 1
    no_dates = sorted(no_hist)
    return {
        "project": project,
        "baseline_date": baseline_date,
        "by_snapshot": {"dates": dates, "created": created, "resolved": resolved,
                        "open": open_cnt, "net": net},
        "by_issue_no": {"dates": no_dates, "created": [no_hist[d] for d in no_dates]},
        "unknown_no": unknown_no,
    }


@router.get("/flow-detail")
def flow_detail(project: str, date: str, kind: str = "created",
                db: Session = Depends(get_db),
                _: models.User = Depends(get_current_user)):
    """某天新增 / 解决的问题单明细。

    新增的单在当天的快照里；解决的单当天已经没了，得回上一次快照的文件里取。
    """
    if kind not in ("created", "resolved"):
        raise HTTPException(400, "kind 只能是 created 或 resolved")
    flow = (
        db.query(models.IssueSnapshotFlow)
        .filter(models.IssueSnapshotFlow.project == project,
                models.IssueSnapshotFlow.snapshot_date == date)
        .first()
    )
    if flow is None:
        return {"project": project, "date": date, "kind": kind, "rows": []}
    ids = set(json.loads((flow.created_ids_json if kind == "created" else flow.resolved_ids_json) or "[]"))
    src_date = date if kind == "created" else (flow.prev_date or "")
    rows: List[Dict] = []
    if ids and src_date:
        snap = (
            db.query(models.IssueSnapshot)
            .filter(models.IssueSnapshot.project == project,
                    models.IssueSnapshot.snapshot_date == src_date)
            .first()
        )
        if snap is not None:
            try:
                fp = _snapshot_root() / snap.data_file
                if fp.exists():
                    rows = [r for r in json.loads(fp.read_text(encoding="utf-8"))
                            if isinstance(r, dict) and str(r.get("issue_id") or "").strip() in ids]
            except Exception:
                rows = []
    return {"project": project, "date": date, "kind": kind,
            "source_date": src_date, "count": len(rows), "rows": rows}


@router.get("/ungrouped")
def ungrouped_owners(project: str, date: Optional[str] = None,
                     db: Session = Depends(get_db),
                     _: models.User = Depends(get_current_user)):
    """某次快照里归不到小组的责任人（默认最新一次），给「小组配置」补名单用。

    从明细文件现算而不是查库：小组名单一改，这份清单就该跟着变，存下来的那份
    会一直显示已经补过的人，比没有更糟。
    """
    q = db.query(models.IssueSnapshot).filter(models.IssueSnapshot.project == project)
    snap = (q.filter(models.IssueSnapshot.snapshot_date == date).first() if date
            else q.order_by(models.IssueSnapshot.snapshot_date.desc()).first())
    if snap is None:
        return {"project": project, "date": "", "rows": [], "count": 0, "issues": 0}
    rows: List[Dict] = []
    try:
        fp = _snapshot_root() / (snap.data_file or "")
        if snap.data_file and fp.exists():
            raw = json.loads(fp.read_text(encoding="utf-8"))
            if isinstance(raw, list):
                rows = _ungrouped_owners([r for r in raw if isinstance(r, dict)])
    except Exception:
        rows = []   # 明细文件丢了：报"没有"而不是整页 500，配置页照常能用
    return {"project": project, "date": snap.snapshot_date, "rows": rows,
            "count": len(rows), "issues": sum(r["count"] for r in rows)}


@router.get("/snapshot-export")
def snapshot_export(request: Request, project: str, date: Optional[str] = None,
                    db: Session = Depends(get_db),
                    current_user: models.User = Depends(get_current_user)):
    """把某次快照导出为 Excel：Sheet1「原始数据」+ Sheet2「统计分析」
    （按小组 / 客户面 / 研发问题 / 年月 × 严重程度）。"""
    import openpyxl

    q = db.query(models.IssueSnapshot).filter(models.IssueSnapshot.project == project)
    snap = (q.filter(models.IssueSnapshot.snapshot_date == date).first() if date
            else q.order_by(models.IssueSnapshot.snapshot_date.desc()).first())
    if snap is None:
        raise HTTPException(404, "该项目暂无快照可导出")
    raw: List[Dict] = []
    try:
        fp = _snapshot_root() / snap.data_file
        if fp.exists():
            raw = json.loads(fp.read_text(encoding="utf-8"))
    except Exception:
        raw = []

    wb = openpyxl.Workbook()
    _fill_raw_sheet(wb.active, raw)
    _fill_analysis_sheet(wb.create_sheet(), raw)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    filename = f"issues_{_safe_slug(project)}_{snap.snapshot_date}.xlsx"
    log_op(db, action="导出Excel", target="问题单", target_id=snap.id,
           detail=f"project={project} date={snap.snapshot_date} rows={len(raw)}",
           user=current_user, request=request)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/run-script/status")
def run_script_status(_: models.User = Depends(get_current_user)):
    """查询刷新脚本是否正在执行（所有登录用户可查）。"""
    return {
        "running":    _script_lock.locked(),
        "started_at": iso_local(_script_started_at),
    }


@router.post("/run-script")
def run_script(
    request: Request,
    db: Session = Depends(get_db),
    current_admin: models.User = Depends(require_admin),
):
    """执行管理员配置的外部刷新脚本（全局互斥，同时只能有一个实例）。"""
    global _script_started_at

    if not _script_lock.acquire(blocking=False):
        raise HTTPException(423, "脚本正在执行中，请等待完成后再试")

    _script_started_at = datetime.utcnow()
    try:
        cfg = _load_config()
        script = cfg.get("issue_script_path", "").strip()
        if not script:
            raise HTTPException(400, "未配置刷新脚本路径（issue_script_path）")

        sp = pathlib.Path(script)
        if not sp.exists():
            raise HTTPException(404, f"脚本不存在：{script}")

        cmd = [sys.executable, str(sp)] if sp.suffix.lower() == ".py" else [str(sp)]

        result = subprocess.run(
            cmd, capture_output=True, text=True,
            timeout=300, cwd=str(sp.parent),
        )
        log_op(db, action="运行脚本", target="问题单",
               detail=f"script={sp.name} exit={result.returncode}",
               user=current_admin, request=request)
        return {
            "ok":        result.returncode == 0,
            "exit_code": result.returncode,
            "stdout":    result.stdout[-3000:] if result.stdout else "",
            "stderr":    result.stderr[-1000:] if result.stderr else "",
        }
    except HTTPException:
        raise
    except subprocess.TimeoutExpired:
        raise HTTPException(500, "脚本执行超时（>5分钟）")
    except Exception as exc:
        raise HTTPException(500, f"脚本启动失败：{exc}")
    finally:
        _script_started_at = None
        _script_lock.release()


@router.get("/export.pptx")
def export_pptx(
    request: Request,
    date: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """将指定日期报表导出为 PPT，不传 date 则取最新。"""
    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        raise HTTPException(400, "未配置报表路径")

    target = _resolve_for_date(path_str, date)
    data   = dict(_parse_excel_cached(str(target)))
    data["actual_file"] = target.name

    try:
        buf = _build_pptx(data)
    except Exception as exc:
        raise HTTPException(500, f"PPT 生成失败：{exc}")
    # HTTP 头是 latin-1 编码，中文文件名必须走 RFC5987（ASCII 回退 + filename*）
    from urllib.parse import quote as url_quote
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"缺陷统计报表_{ts}.pptx"
    log_op(db, action="导出PPT", target="问题单",
           detail=f"date={date or '最新'} file={target.name}",
           user=current_user, request=request)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": (
            f'attachment; filename="issue-report_{ts}.pptx"; '
            f"filename*=UTF-8''{url_quote(filename)}"
        )},
    )
